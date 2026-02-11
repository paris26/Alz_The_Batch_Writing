"""
Build 32-slide editorial-style presentation for AI Alzheimer Thesis.
Uses python-pptx to create a .pptx file with the magazine-style design
specified in PLAN.md.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
BASE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(BASE)
IMG  = os.path.join(REPO, "images")
CHARTS = os.path.join(BASE, "generated_charts")
OUTPUT = os.path.join(BASE, "AI_Alzheimer_Thesis_Presentation.pptx")

# ── Colors ──
DARK_BG    = RGBColor(0x0D, 0x11, 0x17)
LIGHT_BG   = RGBColor(0xF7, 0xF5, 0xF0)
COPPER     = RGBColor(0xC1, 0x7F, 0x3A)
BLUE       = RGBColor(0x3B, 0x82, 0xB6)
RED        = RGBColor(0xDC, 0x4A, 0x4A)
GREEN      = RGBColor(0x5B, 0x8C, 0x6B)
TEXT_ON_DARK  = RGBColor(0xE8, 0xE4, 0xDE)
TEXT_ON_LIGHT = RGBColor(0x2D, 0x2D, 0x2D)
CITE_GRAY  = RGBColor(0x8B, 0x86, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CITE_BAR_DARK  = RGBColor(0x0A, 0x0D, 0x12)
CITE_BAR_LIGHT = RGBColor(0xED, 0xEB, 0xE6)

# ── Dimensions (16:9 widescreen) ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ──

def img_path(filename):
    return os.path.join(IMG, filename)

def chart_path(filename):
    return os.path.join(CHARTS, filename)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_fill(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=TEXT_ON_LIGHT, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, font_name='Georgia',
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    try:
        txBox.text_frame.paragraphs[0].alignment = alignment
    except:
        pass
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    try:
        txBox.text_frame.paragraphs[0].space_before = Pt(0)
        txBox.text_frame.paragraphs[0].space_after = Pt(0)
    except:
        pass
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=14, color=TEXT_ON_LIGHT, bold=False,
                          italic=False, alignment=PP_ALIGN.LEFT,
                          font_name='Georgia', line_spacing=1.2):
    """Add textbox with multiple paragraphs from a list of strings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = font_name
        p.alignment = alignment
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return txBox

def add_citation_strip(slide, text, dark_bg=False):
    """Add architectural citation bar at the bottom of each slide."""
    bar_color = CITE_BAR_DARK if dark_bg else CITE_BAR_LIGHT
    bar = add_shape_fill(slide, Inches(0), SLIDE_H - Inches(0.55),
                         SLIDE_W, Inches(0.55), bar_color)
    bar.line.fill.background()
    add_textbox(slide, Inches(0.5), SLIDE_H - Inches(0.45),
                SLIDE_W - Inches(1), Inches(0.35),
                text, font_size=9, color=CITE_GRAY,
                italic=True, alignment=PP_ALIGN.LEFT)

def add_accent_line(slide, left, top, width, color=COPPER, height=Pt(3)):
    """Thin colored accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_section_label(slide, left, top, text, color=COPPER):
    """Small section label in bold caps."""
    add_textbox(slide, left, top, Inches(3), Inches(0.3),
                text.upper(), font_size=10, color=color, bold=True)

def safe_add_image(slide, img_file, left, top, width=None, height=None):
    """Add image if it exists, skip gracefully if not."""
    path = img_file if os.path.isabs(img_file) else img_path(img_file)
    if not os.path.exists(path):
        print(f"  WARNING: Image not found: {path}")
        return None
    kwargs = {'image_file': path, 'left': left, 'top': top}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height
    return slide.shapes.add_picture(**kwargs)


# ═══════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ──────────────────────────────────────────────
    # SLIDE 1: Title Slide (Dark, Custom)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    # Logo
    safe_add_image(slide, "logo_en.png", Inches(0.6), Inches(0.5),
                   height=Inches(1.0))

    # Accent line
    add_accent_line(slide, Inches(0.6), Inches(2.0), Inches(2.5))

    # Title
    add_textbox(slide, Inches(0.6), Inches(2.3), Inches(10), Inches(1.5),
                "AI Alzheimer and Dementia\nClassification",
                font_size=36, color=WHITE, bold=True)

    # Subtitle
    add_textbox(slide, Inches(0.6), Inches(4.0), Inches(10), Inches(0.8),
                "A Review of Neuroimaging, Machine Learning, and the Road to Clinical Translation",
                font_size=16, color=TEXT_ON_DARK, italic=True)

    # Author block
    add_multiline_textbox(slide, Inches(0.6), Inches(5.2), Inches(8), Inches(1.5),
                          ["Gavriilidis Paraskevas",
                           "Department of Informatics and Computer Engineering",
                           "University of West Attica, 2025"],
                          font_size=14, color=CITE_GRAY)

    add_citation_strip(slide, "Thesis Defense Presentation  •  University of West Attica", dark_bg=True)
    print("  Slide  1: Title")

    # ──────────────────────────────────────────────
    # SLIDE 2: "The Silent Epidemic" (Hero Stat, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "ACT I: THE PROBLEM")

    # Hero number
    add_textbox(slide, Inches(0.6), Inches(1.2), Inches(5), Inches(2.5),
                "7.2M", font_size=84, color=COPPER, bold=True)

    add_textbox(slide, Inches(0.6), Inches(3.5), Inches(5), Inches(0.6),
                "Americans living with Alzheimer's disease",
                font_size=20, color=TEXT_ON_LIGHT, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(4.2), Inches(4))

    add_multiline_textbox(slide, Inches(0.6), Inches(4.5), Inches(5), Inches(1.8),
                          ["Projected to reach 13.8 million by 2060.",
                           "Annual cost exceeds $360 billion.",
                           "Every 65 seconds, someone develops AD."],
                          font_size=14, color=TEXT_ON_LIGHT)

    # Chart on the right
    safe_add_image(slide, chart_path("prevalence_projection.png"),
                   Inches(6.5), Inches(1.0), width=Inches(6.3))

    add_citation_strip(slide,
        "Alzheimer's Association, 2024 Facts and Figures  •  Brookmeyer et al., 2007  •  WHO Dementia Fact Sheet, 2023")
    print("  Slide  2: The Silent Epidemic")

    # ──────────────────────────────────────────────
    # SLIDE 3: "The 20-Year Window" (Statement, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(2),
                "Pathology begins 15–20 years\nbefore the first symptom.",
                font_size=36, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(4.5), Inches(2.3),
                    color=COPPER)

    add_textbox(slide, Inches(2.5), Inches(4.9), Inches(8.3), Inches(1.2),
                "This silent window is both the tragedy and the opportunity — early detection\n"
                "through neuroimaging could transform Alzheimer's from a death sentence to a manageable condition.",
                font_size=15, color=TEXT_ON_DARK, italic=True,
                alignment=PP_ALIGN.CENTER)

    add_citation_strip(slide,
        "Jack et al., 2010, Lancet Neurology  •  Sperling et al., 2011, Alzheimer's & Dementia  •  Bateman et al., 2012, NEJM",
        dark_bg=True)
    print("  Slide  3: The 20-Year Window")

    # ──────────────────────────────────────────────
    # SLIDE 4: "Biomarker Criteria" (Hero variant, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "DIAGNOSTIC FRAMEWORK")

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
                "The AT(N) Framework", font_size=32, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.85), Inches(3))

    # Three columns for A, T, N
    col_w = Inches(3.5)
    col_start = Inches(0.8)
    col_gap = Inches(0.5)

    for i, (letter, title, desc) in enumerate([
        ("A", "Amyloid", "Aβ42 in CSF, amyloid PET — the initiating pathology that seeds decades before symptoms."),
        ("T", "Tau", "Phosphorylated tau in CSF, tau PET — the spreading pathology that tracks with cognitive decline."),
        ("(N)", "Neurodegeneration", "MRI atrophy, FDG-PET hypometabolism — the downstream damage visible on structural imaging.")
    ]):
        x = col_start + i * (col_w + col_gap)
        # Big letter
        add_textbox(slide, x, Inches(2.2), Inches(1.5), Inches(1.5),
                    letter, font_size=60, color=COPPER, bold=True)
        # Title
        add_textbox(slide, x, Inches(3.5), col_w, Inches(0.5),
                    title, font_size=18, color=TEXT_ON_LIGHT, bold=True)
        # Description
        add_textbox(slide, x, Inches(4.1), col_w, Inches(2),
                    desc, font_size=13, color=TEXT_ON_LIGHT)

    # Vertical dividers
    for i in range(1, 3):
        x = col_start + i * (col_w + col_gap) - col_gap / 2
        add_shape_fill(slide, x, Inches(2.4), Pt(1), Inches(3.8), CITE_GRAY)

    add_citation_strip(slide,
        "Jack et al., 2018, Alzheimer's & Dementia (NIA-AA Research Framework)  •  Dubois et al., 2014, Lancet Neurology")
    print("  Slide  4: Biomarker Criteria AT(N)")

    # ──────────────────────────────────────────────
    # SLIDE 5: "Why Neuroimaging" (Full-Bleed, Image-driven)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    # Full-bleed image on left ~60%
    safe_add_image(slide, "nihms-137059-f0004.jpg",
                   Inches(0), Inches(0), width=Inches(8), height=SLIDE_H)

    # Semi-transparent overlay for text on right
    overlay = add_shape_fill(slide, Inches(7), Inches(0), Inches(6.333), SLIDE_H, DARK_BG)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = DARK_BG

    add_section_label(slide, Inches(7.5), Inches(1.0), "THE CASE FOR IMAGING", color=COPPER)

    add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.8),
                "Why Neuroimaging?", font_size=30, color=WHITE, bold=True)

    add_accent_line(slide, Inches(7.5), Inches(2.3), Inches(2))

    add_multiline_textbox(slide, Inches(7.5), Inches(2.7), Inches(5), Inches(3.5),
                          ["Non-invasive — no lumbar punctures, no radioactive tracers for structural MRI.",
                           "",
                           "Quantifiable — voxel-level measurements enable computational analysis at scale.",
                           "",
                           "Objective — reduces inter-rater variability inherent in clinical assessments.",
                           "",
                           "Accessible — MRI scanners exist in most regional hospitals worldwide."],
                          font_size=14, color=TEXT_ON_DARK)

    add_citation_strip(slide,
        "Defined image from Defined et al. NIHMS  •  Defined et al., 2015, Alzheimer's & Dementia  •  Jack et al., 2008, Brain",
        dark_bg=True)
    print("  Slide  5: Why Neuroimaging")

    # ──────────────────────────────────────────────
    # SLIDE 6: "The Datasets" (Split 3-col, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "DATA FOUNDATIONS")

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8),
                "The Benchmark Datasets", font_size=30, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.8), Inches(3))

    datasets = [
        ("ADNI", "Alzheimer's Disease\nNeuroimaging Initiative",
         "2,400+ subjects across 4 phases (2004–present). Multi-site, longitudinal MRI, PET, genetics, CSF biomarkers. The gold standard for AD research.",
         "adni.loni.usc.edu"),
        ("AIBL", "Australian Imaging,\nBiomarker & Lifestyle",
         "1,100+ participants from Melbourne and Perth. Enriched for pre-clinical AD with extensive lifestyle data. Independent validation cohort.",
         "aibl.csiro.au"),
        ("OASIS", "Open Access Series of\nImaging Studies",
         "2,000+ sessions across OASIS-1/2/3/4. Cross-sectional and longitudinal, freely available. Widely used for reproducibility benchmarks.",
         "oasis-brains.org"),
    ]

    col_w = Inches(3.8)
    for i, (name, full, desc, url) in enumerate(datasets):
        x = Inches(0.6) + i * (col_w + Inches(0.25))

        add_textbox(slide, x, Inches(2.2), col_w, Inches(0.7),
                    name, font_size=36, color=COPPER, bold=True)
        add_textbox(slide, x, Inches(2.9), col_w, Inches(0.6),
                    full, font_size=12, color=CITE_GRAY, italic=True)
        add_accent_line(slide, x, Inches(3.55), Inches(1.5))
        add_textbox(slide, x, Inches(3.7), col_w, Inches(2.5),
                    desc, font_size=13, color=TEXT_ON_LIGHT)
        add_textbox(slide, x, Inches(5.8), col_w, Inches(0.3),
                    url, font_size=10, color=BLUE, italic=True)

    # Vertical dividers
    for i in range(1, 3):
        x = Inches(0.6) + i * (col_w + Inches(0.25)) - Inches(0.12)
        add_shape_fill(slide, x, Inches(2.2), Pt(1), Inches(4), CITE_GRAY)

    add_citation_strip(slide,
        "Mueller et al., 2005, Neuroimaging Clin N Am (ADNI)  •  Ellis et al., 2009, Int Psychogeriatr (AIBL)  •  Marcus et al., 2007, J Cogn Neurosci (OASIS)")
    print("  Slide  6: The Datasets")

    # ──────────────────────────────────────────────
    # SLIDE 7: Section Divider — "SEEING THE BRAIN" (Statement, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1), Inches(0.6), Inches(4), Inches(0.4),
                "ACT II", font_size=12, color=COPPER, bold=True)

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(2),
                "SEEING THE BRAIN", font_size=52, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(4.6), Inches(2.3))

    add_textbox(slide, Inches(2.5), Inches(5.0), Inches(8.3), Inches(1),
                "Imaging modalities, preprocessing, and the computational pipeline",
                font_size=16, color=TEXT_ON_DARK, italic=True,
                alignment=PP_ALIGN.CENTER)

    add_citation_strip(slide, "Section 2: Neuroimaging Modalities and Preprocessing", dark_bg=True)
    print("  Slide  7: Section Divider - SEEING THE BRAIN")

    # ──────────────────────────────────────────────
    # SLIDE 8: "MRI Fundamentals" (Full-Bleed, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    safe_add_image(slide, "IntensityNormalization1.png",
                   Inches(0), Inches(0), width=Inches(7.5), height=SLIDE_H)

    overlay = add_shape_fill(slide, Inches(6.8), Inches(0), Inches(6.533), SLIDE_H, DARK_BG)

    add_section_label(slide, Inches(7.3), Inches(0.8), "IMAGING MODALITIES", color=COPPER)

    add_textbox(slide, Inches(7.3), Inches(1.3), Inches(5.5), Inches(0.8),
                "MRI Fundamentals", font_size=30, color=WHITE, bold=True)

    add_accent_line(slide, Inches(7.3), Inches(2.1), Inches(2))

    add_multiline_textbox(slide, Inches(7.3), Inches(2.5), Inches(5.3), Inches(4),
                          ["Magnetic Resonance Imaging exploits the Larmor equation: "
                           "nuclei precess at frequencies proportional to field strength.",
                           "",
                           "T1-weighted scans provide excellent gray/white matter contrast — "
                           "the workhorse of structural neuroimaging.",
                           "",
                           "T2-weighted and FLAIR sequences detect pathological fluid accumulation "
                           "and white matter lesions.",
                           "",
                           "Modern 3T scanners achieve sub-millimeter resolution, enabling "
                           "voxel-level analysis of atrophy patterns."],
                          font_size=13, color=TEXT_ON_DARK)

    add_citation_strip(slide,
        "Defined, 2002, MRI: Basic Principles  •  Jack et al., 2008, Brain  •  Defined, 2015, Neuroimaging",
        dark_bg=True)
    print("  Slide  8: MRI Fundamentals")

    # ──────────────────────────────────────────────
    # SLIDE 9: "CT and PET" (Split Compare, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "MODALITY COMPARISON")

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
                "CT vs PET Imaging", font_size=30, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(3))

    # Left: CT
    safe_add_image(slide, "pmp-32-1-1-f1.png",
                   Inches(0.6), Inches(2.2), width=Inches(5.5))

    add_textbox(slide, Inches(0.6), Inches(5.6), Inches(5.5), Inches(0.4),
                "CT — Structural Anatomy", font_size=16, color=TEXT_ON_LIGHT, bold=True)
    add_textbox(slide, Inches(0.6), Inches(6.0), Inches(5.5), Inches(0.7),
                "Fast, widely available. Shows calcifications and acute hemorrhage. Limited soft-tissue contrast makes it secondary to MRI for dementia assessment.",
                font_size=12, color=TEXT_ON_LIGHT)

    # Divider
    add_shape_fill(slide, Inches(6.5), Inches(2.2), Pt(1.5), Inches(4.5), CITE_GRAY)

    # Right: PET
    safe_add_image(slide, "pmp-32-1-1-f4.png",
                   Inches(7), Inches(2.2), width=Inches(5.5))

    add_textbox(slide, Inches(7), Inches(5.6), Inches(5.5), Inches(0.4),
                "PET — Molecular Function", font_size=16, color=TEXT_ON_LIGHT, bold=True)
    add_textbox(slide, Inches(7), Inches(6.0), Inches(5.5), Inches(0.7),
                "Reveals metabolic activity (FDG) or protein deposits (amyloid/tau tracers). Essential for early detection but expensive and requires radioactive tracers.",
                font_size=12, color=TEXT_ON_LIGHT)

    add_citation_strip(slide,
        "Park et al., 2020, Prog Med Phys  •  Johnson et al., 2012, Ann Neurol  •  Defined, 2017, Brain Imaging")
    print("  Slide  9: CT and PET")

    # ──────────────────────────────────────────────
    # SLIDE 10: "PET Biomarkers" (Hero variant, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "MOLECULAR IMAGING")

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
                "PET Biomarkers", font_size=30, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(2.5))

    # Hero number
    add_textbox(slide, Inches(0.6), Inches(2.3), Inches(5), Inches(2),
                "90%", font_size=80, color=COPPER, bold=True)

    add_textbox(slide, Inches(0.6), Inches(4.2), Inches(5), Inches(0.5),
                "FDG-PET sensitivity for AD detection",
                font_size=18, color=TEXT_ON_LIGHT, bold=True)

    # Right column with three biomarker types
    biomarkers = [
        ("FDG-PET", "Measures glucose metabolism. Hypometabolism in temporo-parietal regions is a hallmark of AD. Sensitivity ~90%, specificity ~71%."),
        ("Amyloid PET", "Pittsburgh Compound B (PiB), Florbetapir. Detects amyloid plaques 15–20 years before symptoms. Changed clinical trial enrollment."),
        ("Tau PET", "Flortaucipir (AV-1451). Maps neurofibrillary tangle distribution. Correlates more closely with cognitive decline than amyloid."),
    ]

    for i, (name, desc) in enumerate(biomarkers):
        y = Inches(2.2) + i * Inches(1.6)
        add_textbox(slide, Inches(6.5), y, Inches(6), Inches(0.4),
                    name, font_size=16, color=COPPER, bold=True)
        add_textbox(slide, Inches(6.5), y + Inches(0.4), Inches(6), Inches(1),
                    desc, font_size=12, color=TEXT_ON_LIGHT)

    add_citation_strip(slide,
        "Defined et al., 2004, J Nucl Med (FDG)  •  Klunk et al., 2004, Ann Neurol (PiB)  •  Johnson et al., 2016, Ann Neurol (Tau)")
    print("  Slide 10: PET Biomarkers")

    # ──────────────────────────────────────────────
    # SLIDE 11: "Preprocessing Pipeline" (Dark Canvas)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_section_label(slide, Inches(0.6), Inches(0.5), "PREPROCESSING", color=COPPER)

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
                "The Preprocessing Pipeline", font_size=30, color=WHITE, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.8), Inches(2.5))

    # Pipeline steps as a visual flow
    steps = [
        ("1", "Raw MRI\nAcquisition", BLUE),
        ("→", "", CITE_GRAY),
        ("2", "Intensity\nNormalization", COPPER),
        ("→", "", CITE_GRAY),
        ("3", "Denoising\n(NLM/BM3D)", GREEN),
        ("→", "", CITE_GRAY),
        ("4", "Skull\nStripping", RED),
        ("→", "", CITE_GRAY),
        ("5", "Registration &\nSegmentation", BLUE),
    ]

    x_start = Inches(0.5)
    y_center = Inches(3.2)

    for i, (num, label, color) in enumerate(steps):
        x = x_start + i * Inches(1.38)
        if num == "→":
            add_textbox(slide, x, y_center + Inches(0.2), Inches(0.8), Inches(0.6),
                        "→", font_size=30, color=CITE_GRAY, bold=True,
                        alignment=PP_ALIGN.CENTER)
        else:
            # Box
            box = add_shape_fill(slide, x, y_center, Inches(1.2), Inches(1.4), RGBColor(0x16, 0x1B, 0x22))
            box.line.color.rgb = color
            box.line.width = Pt(2)
            # Number
            add_textbox(slide, x, y_center + Inches(0.1), Inches(1.2), Inches(0.5),
                        num, font_size=24, color=color, bold=True,
                        alignment=PP_ALIGN.CENTER)
            # Label
            add_textbox(slide, x, y_center + Inches(0.5), Inches(1.2), Inches(0.8),
                        label, font_size=11, color=TEXT_ON_DARK,
                        alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(1),
                "Each step introduces assumptions and potential artifacts. Inconsistent preprocessing is a leading cause "
                "of irreproducible results in neuroimaging ML studies.",
                font_size=14, color=TEXT_ON_DARK, italic=True)

    add_citation_strip(slide,
        "Defined et al., 2019, NeuroImage  •  Defined et al., 2014, Frontiers in Neuroscience  •  Ashburner, 2007, NeuroImage",
        dark_bg=True)
    print("  Slide 11: Preprocessing Pipeline")

    # ──────────────────────────────────────────────
    # SLIDE 12: "Intensity Normalization" (Split Compare, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "PREPROCESSING")

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
                "Intensity Normalization", font_size=30, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(3))

    # Left image
    safe_add_image(slide, "IntensityNormalization2.png",
                   Inches(0.6), Inches(2.2), width=Inches(5.8))

    # Divider
    add_shape_fill(slide, Inches(6.55), Inches(2.2), Pt(1.5), Inches(4.3), CITE_GRAY)

    # Right image
    safe_add_image(slide, "IntensityNormalization3.png",
                   Inches(7), Inches(2.2), width=Inches(5.8))

    # Labels below
    methods = [
        ("Z-Score Normalization", "Subtracts mean, divides by standard deviation. Simple but assumes Gaussian intensity distribution."),
        ("Histogram Matching", "Aligns intensity distributions across scans. Robust to scanner variability but can distort pathological signals."),
        ("White Stripe", "Uses normal-appearing white matter as internal reference. Most principled approach for multi-site studies."),
    ]

    for i, (name, desc) in enumerate(methods):
        x = Inches(0.6) + i * Inches(4.2)
        add_textbox(slide, x, Inches(5.7), Inches(3.8), Inches(0.3),
                    name, font_size=12, color=COPPER, bold=True)
        add_textbox(slide, x, Inches(6.0), Inches(3.8), Inches(0.7),
                    desc, font_size=10, color=TEXT_ON_LIGHT)

    add_citation_strip(slide,
        "Shinohara et al., 2014, NeuroImage (White Stripe)  •  Nyúl et al., 2000, IEEE TMI  •  Shah et al., 2011, J Neurosci Methods")
    print("  Slide 12: Intensity Normalization")

    # ──────────────────────────────────────────────
    # SLIDE 13: "Denoising" (Hero variant, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "PREPROCESSING")

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
                "Denoising Strategies", font_size=30, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(2.5))

    # Three denoising methods
    methods = [
        ("NLM", "Non-Local Means",
         "Exploits self-similarity across patches. Weighted average of similar neighborhoods. "
         "Preserves edges better than Gaussian filtering. O(n²) complexity."),
        ("BM3D", "Block-Matching 3D",
         "Groups similar 2D blocks into 3D arrays, applies collaborative filtering in transform domain. "
         "State-of-the-art for natural images, adapted for MRI."),
        ("DL", "Deep Learning",
         "Encoder-decoder networks (DnCNN, U-Net variants). Learn noise patterns from paired data. "
         "Fastest inference but require training data matching target scanner."),
    ]

    for i, (abbrev, full, desc) in enumerate(methods):
        x = Inches(0.6) + i * Inches(4.2)
        add_textbox(slide, x, Inches(2.2), Inches(1.5), Inches(1),
                    abbrev, font_size=42, color=COPPER, bold=True)
        add_textbox(slide, x, Inches(3.2), Inches(3.8), Inches(0.4),
                    full, font_size=14, color=TEXT_ON_LIGHT, bold=True)
        add_accent_line(slide, x, Inches(3.65), Inches(1.5))
        add_textbox(slide, x, Inches(3.9), Inches(3.8), Inches(2.5),
                    desc, font_size=12, color=TEXT_ON_LIGHT)

    # Vertical dividers
    for i in range(1, 3):
        x = Inches(0.6) + i * Inches(4.2) - Inches(0.1)
        add_shape_fill(slide, x, Inches(2.2), Pt(1), Inches(4), CITE_GRAY)

    add_citation_strip(slide,
        "Buades et al., 2005, CVPR (NLM)  •  Dabov et al., 2007, IEEE TIP (BM3D)  •  Zhang et al., 2017, IEEE TIP (DnCNN)")
    print("  Slide 13: Denoising")

    # ──────────────────────────────────────────────
    # SLIDE 14: "Skull Stripping" (Full-Bleed, Image-driven)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    safe_add_image(slide, "Skull Stripping image.png",
                   Inches(0), Inches(0), width=Inches(7), height=SLIDE_H)

    overlay = add_shape_fill(slide, Inches(6.2), Inches(0), Inches(7.133), SLIDE_H, DARK_BG)

    add_section_label(slide, Inches(6.8), Inches(0.8), "PREPROCESSING", color=COPPER)

    add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.7),
                "Skull Stripping", font_size=30, color=WHITE, bold=True)

    add_accent_line(slide, Inches(6.8), Inches(2.05), Inches(2))

    add_multiline_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.8), Inches(3),
                          ["Removal of non-brain tissue (skull, scalp, dura mater) is critical — "
                           "residual skull can dominate classifier features.",
                           "",
                           "Classical: BET (FSL), FreeSurfer watershed. Robust but slow and require manual QC.",
                           "",
                           "Deep learning: HD-BET, SynthStrip achieve Dice >0.97 with zero manual intervention.",
                           "",
                           "Failure modes: over-stripping removes cortex, under-stripping leaves meninges. "
                           "Both corrupt downstream VBM and classification."],
                          font_size=13, color=TEXT_ON_DARK)

    safe_add_image(slide, "Skull Stripping Techniques.png",
                   Inches(6.8), Inches(5.2), width=Inches(5.5), height=Inches(1.5))

    add_citation_strip(slide,
        "Smith, 2002, HBM (BET)  •  Isensee et al., 2019, NeuroImage (HD-BET)  •  Hoopes et al., 2022, NeuroImage (SynthStrip)",
        dark_bg=True)
    print("  Slide 14: Skull Stripping")

    # ──────────────────────────────────────────────
    # SLIDE 15: "Voxel-Based Morphometry" (Full-Bleed, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    safe_add_image(slide, "nihms154848f1.jpg",
                   Inches(0), Inches(0), width=Inches(7.5), height=SLIDE_H)

    overlay = add_shape_fill(slide, Inches(6.8), Inches(0), Inches(6.533), SLIDE_H, DARK_BG)

    add_section_label(slide, Inches(7.3), Inches(0.8), "FEATURE EXTRACTION", color=COPPER)

    add_textbox(slide, Inches(7.3), Inches(1.3), Inches(5.5), Inches(0.7),
                "Voxel-Based Morphometry", font_size=28, color=WHITE, bold=True)

    add_accent_line(slide, Inches(7.3), Inches(2.05), Inches(2))

    add_multiline_textbox(slide, Inches(7.3), Inches(2.5), Inches(5.3), Inches(4),
                          ["VBM quantifies regional gray matter concentration differences "
                           "across the entire brain, without pre-selecting regions of interest.",
                           "",
                           "Pipeline: Segmentation → DARTEL normalization → Smoothing → "
                           "Statistical parametric mapping.",
                           "",
                           "AD signature: bilateral hippocampal atrophy, entorhinal cortex thinning, "
                           "temporo-parietal gray matter loss.",
                           "",
                           "AUC >0.90 for AD vs HC classification using VBM features alone — "
                           "validating structural imaging as a powerful biomarker."],
                          font_size=13, color=TEXT_ON_DARK)

    add_citation_strip(slide,
        "Ashburner & Friston, 2000, NeuroImage (VBM)  •  Karas et al., 2004, NeuroImage  •  Defined, 2018, Alzheimer's Research & Therapy",
        dark_bg=True)
    print("  Slide 15: Voxel-Based Morphometry")

    # ──────────────────────────────────────────────
    # SLIDE 16: "Classical ML" (Hero Stat, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "MACHINE LEARNING")

    # Hero number
    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(5), Inches(2.5),
                "94.5%", font_size=80, color=COPPER, bold=True)

    add_textbox(slide, Inches(0.6), Inches(3.2), Inches(5), Inches(0.5),
                "SVM accuracy for AD vs Healthy Controls",
                font_size=18, color=TEXT_ON_LIGHT, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(3.8), Inches(3.5))

    add_multiline_textbox(slide, Inches(0.6), Inches(4.1), Inches(5.5), Inches(2),
                          ["But the MCI cliff tells the real story — accuracy drops "
                           "to ~68% for mild cognitive impairment detection.",
                           "",
                           "SVM, Random Forest, and Logistic Regression all show the "
                           "same pattern: binary AD detection is solved; the clinically "
                           "relevant task (early MCI) remains elusive."],
                          font_size=13, color=TEXT_ON_LIGHT)

    # Chart on the right
    safe_add_image(slide, chart_path("ml_comparison.png"),
                   Inches(6.5), Inches(0.8), width=Inches(6.3))

    add_citation_strip(slide,
        "Defined et al., 2014, NeuroImage  •  Defined et al., 2018, IEEE  •  Defined et al., 2017, Alzheimer's & Dementia")
    print("  Slide 16: Classical ML")

    # ──────────────────────────────────────────────
    # SLIDE 17: "Deep Learning Revolution" (Dark Canvas)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_section_label(slide, Inches(0.6), Inches(0.5), "DEEP LEARNING", color=COPPER)

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
                "The Deep Learning Revolution", font_size=30, color=WHITE, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.8), Inches(2.5))

    # Architecture comparison boxes
    archs = [
        ("2D CNN", "Slice-level", "Process individual 2D slices. Fast training, large effective dataset "
         "(N × slices). Miss inter-slice spatial relationships. Risk of information leakage between slices of same subject."),
        ("3D CNN", "Volume-level", "Process entire 3D brain volumes. Capture spatial context across all axes. "
         "Require more GPU memory and larger datasets. More prone to overfitting on small cohorts."),
        ("ResNet/DenseNet", "Transfer Learning", "Pre-trained on ImageNet, fine-tuned on brain scans. "
         "Exploit low-level feature reuse. Dominant approach: 2D pretrained → slice-level classification → majority voting."),
    ]

    for i, (name, subtitle, desc) in enumerate(archs):
        x = Inches(0.5) + i * Inches(4.2)
        box = add_shape_fill(slide, x, Inches(2.4), Inches(3.9), Inches(3.8),
                             RGBColor(0x16, 0x1B, 0x22))
        box.line.color.rgb = BLUE if i < 2 else GREEN
        box.line.width = Pt(1.5)

        add_textbox(slide, x + Inches(0.2), Inches(2.6), Inches(3.5), Inches(0.5),
                    name, font_size=22, color=BLUE if i < 2 else GREEN, bold=True)
        add_textbox(slide, x + Inches(0.2), Inches(3.1), Inches(3.5), Inches(0.3),
                    subtitle, font_size=12, color=COPPER, italic=True)
        add_accent_line(slide, x + Inches(0.2), Inches(3.45), Inches(1.5),
                        color=BLUE if i < 2 else GREEN)
        add_textbox(slide, x + Inches(0.2), Inches(3.7), Inches(3.5), Inches(2.2),
                    desc, font_size=12, color=TEXT_ON_DARK)

    add_citation_strip(slide,
        "He et al., 2016, CVPR (ResNet)  •  Huang et al., 2017, CVPR (DenseNet)  •  Wen et al., 2020, Med Image Anal",
        dark_bg=True)
    print("  Slide 17: Deep Learning Revolution")

    # ──────────────────────────────────────────────
    # SLIDE 18: "Vision Transformers & Hybrids" (Split, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "ADVANCED ARCHITECTURES")

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
                "Vision Transformers & Hybrid Models", font_size=28, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(3))

    # Left: ViT
    add_textbox(slide, Inches(0.6), Inches(2.2), Inches(5.5), Inches(0.5),
                "Vision Transformer (ViT)", font_size=20, color=BLUE, bold=True)

    add_multiline_textbox(slide, Inches(0.6), Inches(2.8), Inches(5.5), Inches(3.5),
                          ["Patches brain images into 16×16 tokens, processes via self-attention.",
                           "",
                           "Global receptive field from layer one — captures long-range "
                           "spatial dependencies that CNNs build gradually.",
                           "",
                           "Data hungry: requires large datasets or pre-training. "
                           "Performance degrades sharply on small cohorts (<500 subjects).",
                           "",
                           "Emerging evidence: attention maps align with known AD atrophy regions."],
                          font_size=13, color=TEXT_ON_LIGHT)

    # Divider
    add_shape_fill(slide, Inches(6.4), Inches(2.2), Pt(1.5), Inches(4.3), CITE_GRAY)

    # Right: Hybrid
    add_textbox(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(0.5),
                "Hybrid CNN + Transformer", font_size=20, color=GREEN, bold=True)

    add_multiline_textbox(slide, Inches(6.8), Inches(2.8), Inches(5.8), Inches(3.5),
                          ["CNN backbone extracts local features; transformer head captures global context.",
                           "",
                           "Best of both worlds: inductive bias from convolutions, "
                           "long-range attention from transformers.",
                           "",
                           "CNN+SVM hybrid remains competitive: CNN as feature extractor, "
                           "SVM as classifier. Simpler, more interpretable.",
                           "",
                           "Trend: hybrid architectures increasingly dominate leaderboards "
                           "for medical imaging tasks."],
                          font_size=13, color=TEXT_ON_LIGHT)

    add_citation_strip(slide,
        "Dosovitskiy et al., 2021, ICLR (ViT)  •  Defined et al., 2022, Med Image Anal  •  Defined et al., 2023, NeuroImage")
    print("  Slide 18: Vision Transformers & Hybrids")

    # ──────────────────────────────────────────────
    # SLIDE 19: "Measuring Performance" (Hero variant, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "EVALUATION")

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(6), Inches(0.6),
                "Measuring Performance", font_size=30, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(2.5))

    add_multiline_textbox(slide, Inches(0.6), Inches(2.2), Inches(5.5), Inches(4),
                          ["Accuracy alone is misleading when classes are imbalanced — "
                           "a model predicting \"no AD\" achieves 85% accuracy trivially.",
                           "",
                           "Essential metrics for clinical AD classification:",
                           "",
                           "Sensitivity (Recall) — proportion of true AD cases detected. "
                           "Missing a diagnosis has devastating consequences.",
                           "",
                           "Specificity — proportion of healthy controls correctly identified. "
                           "False positives cause unnecessary anxiety and procedures.",
                           "",
                           "AUC-ROC — threshold-independent measure of discriminative power. "
                           "The gold standard for comparing models."],
                          font_size=13, color=TEXT_ON_LIGHT)

    # ROC chart on the right
    safe_add_image(slide, chart_path("roc_curve.png"),
                   Inches(7), Inches(1.0), width=Inches(5.5))

    add_citation_strip(slide,
        "Defined et al., 2019, Lancet Digital Health  •  Defined, 2017, BMC Med Inform  •  Defined et al., 2020, J Alzheimer's Dis")
    print("  Slide 19: Measuring Performance")

    # ──────────────────────────────────────────────
    # SLIDE 20: "Explainability" (Full-Bleed, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    safe_add_image(slide, "Grad-CAMVBM.png",
                   Inches(0), Inches(0), width=Inches(7), height=SLIDE_H)

    overlay = add_shape_fill(slide, Inches(6.2), Inches(0), Inches(7.133), SLIDE_H, DARK_BG)

    add_section_label(slide, Inches(6.8), Inches(0.8), "EXPLAINABILITY", color=COPPER)

    add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.7),
                "Opening the Black Box", font_size=28, color=WHITE, bold=True)

    add_accent_line(slide, Inches(6.8), Inches(2.05), Inches(2))

    xai_methods = [
        ("Grad-CAM", "Gradient-weighted class activation maps highlight which brain regions "
         "drive the classification decision. Visual, intuitive, but coarse resolution."),
        ("LIME", "Locally Interpretable Model-agnostic Explanations perturb input regions and observe "
         "output changes. Model-agnostic but computationally expensive for 3D volumes."),
        ("SHAP", "SHapley Additive exPlanations provide theoretically grounded feature attributions. "
         "Consistent and locally accurate, but prohibitively slow for large models."),
    ]

    for i, (name, desc) in enumerate(xai_methods):
        y = Inches(2.5) + i * Inches(1.5)
        add_textbox(slide, Inches(6.8), y, Inches(6), Inches(0.4),
                    name, font_size=16, color=COPPER, bold=True)
        add_textbox(slide, Inches(6.8), y + Inches(0.35), Inches(5.8), Inches(1),
                    desc, font_size=12, color=TEXT_ON_DARK)

    add_textbox(slide, Inches(6.8), Inches(5.8), Inches(5.8), Inches(0.8),
                "The three pillars: Transparency (how it works), Justification (why this prediction), "
                "and Informativeness (what we learn about the disease).",
                font_size=12, color=TEXT_ON_DARK, italic=True)

    add_citation_strip(slide,
        "Selvaraju et al., 2017, ICCV (Grad-CAM)  •  Ribeiro et al., 2016, KDD (LIME)  •  Lundberg & Lee, 2017, NeurIPS (SHAP)",
        dark_bg=True)
    print("  Slide 20: Explainability")

    # ──────────────────────────────────────────────
    # SLIDE 21: "Our Experiments" (Full-Bleed, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    safe_add_image(slide, "limitations_gradcam.png",
                   Inches(0), Inches(0), width=Inches(7.5), height=SLIDE_H)

    overlay = add_shape_fill(slide, Inches(6.8), Inches(0), Inches(6.533), SLIDE_H, DARK_BG)

    add_section_label(slide, Inches(7.3), Inches(0.8), "EXPERIMENTS", color=COPPER)

    add_textbox(slide, Inches(7.3), Inches(1.3), Inches(5.5), Inches(0.7),
                "Grad-CAM Across\nDementia Stages", font_size=28, color=WHITE, bold=True)

    add_accent_line(slide, Inches(7.3), Inches(2.3), Inches(2))

    add_multiline_textbox(slide, Inches(7.3), Inches(2.7), Inches(5.3), Inches(3.5),
                          ["Our Grad-CAM analysis reveals attention patterns across four "
                           "dementia stages from the OASIS dataset.",
                           "",
                           "AD stage: Model focuses on medial temporal lobe — anatomically "
                           "consistent with known atrophy patterns.",
                           "",
                           "MCI stage: Attention is diffuse and inconsistent — reflecting "
                           "the genuine diagnostic ambiguity of this stage.",
                           "",
                           "Key finding: Models learn stage-specific visual signatures, but "
                           "the clinical utility depends on whether these signatures generalize "
                           "beyond the training distribution."],
                          font_size=13, color=TEXT_ON_DARK)

    add_citation_strip(slide,
        "Experimental results from thesis Chapter 8  •  OASIS-3 dataset  •  ResNet-50 backbone with Grad-CAM visualization",
        dark_bg=True)
    print("  Slide 21: Our Experiments")

    # ──────────────────────────────────────────────
    # SLIDE 22: Section Divider — "THE INCONVENIENT TRUTHS"
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1), Inches(0.6), Inches(4), Inches(0.4),
                "ACT III", font_size=12, color=RED, bold=True)

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(2),
                "THE INCONVENIENT\nTRUTHS", font_size=48, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(4.8), Inches(2.3), color=RED)

    add_textbox(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(1),
                "Methodological pitfalls, data quality, and the gap between benchmarks and bedside",
                font_size=16, color=TEXT_ON_DARK, italic=True,
                alignment=PP_ALIGN.CENTER)

    add_citation_strip(slide, "Section 3: Critical Analysis and Limitations", dark_bg=True)
    print("  Slide 22: Section Divider - THE INCONVENIENT TRUTHS")

    # ──────────────────────────────────────────────
    # SLIDE 23: "Data Leakage" (Hero Stat, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "CRITICAL FLAW", color=RED)

    # Hero number
    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(5), Inches(2.5),
                "−28%", font_size=80, color=RED, bold=True)

    add_textbox(slide, Inches(0.6), Inches(3.2), Inches(5), Inches(0.5),
                "Accuracy drop when data leakage is corrected",
                font_size=18, color=TEXT_ON_LIGHT, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(3.8), Inches(3.5), color=RED)

    add_multiline_textbox(slide, Inches(0.6), Inches(4.1), Inches(5.5), Inches(2.2),
                          ["Only 4.5% of published studies use proper subject-level "
                           "train/test splitting.",
                           "",
                           "When slices from the same patient appear in both training "
                           "and test sets, models memorize patient identity rather than "
                           "learning disease biomarkers.",
                           "",
                           "This single methodological flaw invalidates a majority of "
                           "reported results in the literature."],
                          font_size=13, color=TEXT_ON_LIGHT)

    # Chart on the right
    safe_add_image(slide, chart_path("data_leakage_impact.png"),
                   Inches(6.5), Inches(0.8), width=Inches(6.3))

    add_citation_strip(slide,
        "Wen et al., 2020, Med Image Anal  •  Yagis et al., 2021, J Neurosci Methods  •  Defined et al., 2022, NeuroImage")
    print("  Slide 23: Data Leakage")

    # ──────────────────────────────────────────────
    # SLIDE 24: "Accuracy Paradox" (Full-Bleed, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "EVALUATION PITFALL", color=RED)

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(6), Inches(0.6),
                "The Accuracy Paradox", font_size=30, color=RED, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(2.5), color=RED)

    safe_add_image(slide, "limitations_class_imbalance.png",
                   Inches(0.4), Inches(2.2), width=Inches(6.5))

    add_multiline_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5),
                          ["Class imbalance is endemic in AD datasets. The Moderate "
                           "Dementia class represents just ~1% of OASIS subjects.",
                           "",
                           "A naive model predicting 'Nondemented' for every scan achieves "
                           "~85% accuracy — a meaningless number that appears impressive.",
                           "",
                           "The clinically relevant minority classes (MCI, Moderate AD) "
                           "are precisely the ones where models fail most catastrophically.",
                           "",
                           "Solution: Report balanced accuracy, F1-macro, sensitivity per class. "
                           "Never trust a single accuracy number without understanding "
                           "the class distribution."],
                          font_size=13, color=TEXT_ON_LIGHT)

    add_citation_strip(slide,
        "OASIS-3 demographics  •  Defined et al., 2020, Med Image Anal  •  Defined, 2019, Applied Sciences")
    print("  Slide 24: Accuracy Paradox")

    # ──────────────────────────────────────────────
    # SLIDE 25: "Domain Shift" (Hero Stat, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "GENERALIZATION GAP", color=RED)

    add_textbox(slide, Inches(0.6), Inches(1.2), Inches(5), Inches(2.5),
                "71%", font_size=84, color=RED, bold=True)

    add_textbox(slide, Inches(0.6), Inches(3.5), Inches(5), Inches(0.5),
                "Accuracy on external validation data",
                font_size=18, color=TEXT_ON_LIGHT, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(4.1), Inches(3.5), color=RED)

    add_multiline_textbox(slide, Inches(0.6), Inches(4.4), Inches(5.5), Inches(2),
                          ["Models trained on ADNI and tested on independent datasets "
                           "lose 20–30% accuracy.",
                           "",
                           "Scanner differences, acquisition protocols, demographics, "
                           "and preprocessing choices all contribute to domain shift."],
                          font_size=13, color=TEXT_ON_LIGHT)

    # Right side: comparison
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                "The Generalization Problem", font_size=22, color=RED, bold=True)

    add_accent_line(slide, Inches(7), Inches(2.1), Inches(2), color=RED)

    comparisons = [
        ("ADNI → ADNI", "~94% accuracy", "(same distribution)", GREEN),
        ("ADNI → AIBL", "~78% accuracy", "(similar demographics)", COPPER),
        ("ADNI → OASIS", "~71% accuracy", "(different scanners)", RED),
        ("ADNI → Clinical", "Unknown", "(no systematic testing)", CITE_GRAY),
    ]

    for i, (test, acc, note, color) in enumerate(comparisons):
        y = Inches(2.5) + i * Inches(1.1)
        add_textbox(slide, Inches(7), y, Inches(3), Inches(0.4),
                    test, font_size=16, color=TEXT_ON_LIGHT, bold=True)
        add_textbox(slide, Inches(10), y, Inches(2.5), Inches(0.4),
                    acc, font_size=16, color=color, bold=True)
        add_textbox(slide, Inches(7), y + Inches(0.35), Inches(5.5), Inches(0.4),
                    note, font_size=11, color=CITE_GRAY, italic=True)

    add_citation_strip(slide,
        "Defined et al., 2020, Med Image Anal  •  Defined et al., 2019, NeuroImage  •  Defined et al., 2023, Alzheimer's & Dementia")
    print("  Slide 25: Domain Shift")

    # ──────────────────────────────────────────────
    # SLIDE 26: "Shortcut Learning" (Dark Canvas)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_section_label(slide, Inches(0.6), Inches(0.5), "HIDDEN DANGER", color=RED)

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
                "Shortcut Learning", font_size=30, color=WHITE, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.8), Inches(2.5), color=RED)

    # Two comparison boxes
    # Left: Correct reasoning
    box_l = add_shape_fill(slide, Inches(0.5), Inches(2.4), Inches(5.8), Inches(3.5),
                           RGBColor(0x16, 0x1B, 0x22))
    box_l.line.color.rgb = GREEN
    box_l.line.width = Pt(2)

    add_textbox(slide, Inches(0.7), Inches(2.6), Inches(5.4), Inches(0.5),
                "Correct Reasoning", font_size=20, color=GREEN, bold=True)

    add_multiline_textbox(slide, Inches(0.7), Inches(3.2), Inches(5.4), Inches(2.5),
                          ["Model learns hippocampal atrophy patterns",
                           "Focuses on medial temporal lobe",
                           "Attention maps align with clinical knowledge",
                           "Generalizes to new populations"],
                          font_size=13, color=TEXT_ON_DARK)

    # Right: Shortcut reasoning
    box_r = add_shape_fill(slide, Inches(7), Inches(2.4), Inches(5.8), Inches(3.5),
                           RGBColor(0x16, 0x1B, 0x22))
    box_r.line.color.rgb = RED
    box_r.line.width = Pt(2)

    add_textbox(slide, Inches(7.2), Inches(2.6), Inches(5.4), Inches(0.5),
                "Shortcut Reasoning", font_size=20, color=RED, bold=True)

    add_multiline_textbox(slide, Inches(7.2), Inches(3.2), Inches(5.4), Inches(2.5),
                          ["Model learns scanner artifacts or head position",
                           "Exploits correlation between age and AD prevalence",
                           "Attention maps show skull edges or background",
                           "Fails catastrophically on new scanners"],
                          font_size=13, color=TEXT_ON_DARK)

    # VS divider
    add_textbox(slide, Inches(5.9), Inches(3.5), Inches(1.5), Inches(0.8),
                "vs", font_size=28, color=COPPER, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.7),
                "Without rigorous explainability analysis, we cannot distinguish these two scenarios from accuracy alone.",
                font_size=14, color=TEXT_ON_DARK, italic=True,
                alignment=PP_ALIGN.CENTER)

    add_citation_strip(slide,
        "Geirhos et al., 2020, Nature Machine Intelligence  •  DeGrave et al., 2021, Nature Machine Intelligence",
        dark_bg=True)
    print("  Slide 26: Shortcut Learning")

    # ──────────────────────────────────────────────
    # SLIDE 27: "Label Uncertainty" (Hero Stat, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "GROUND TRUTH PROBLEM", color=RED)

    add_textbox(slide, Inches(0.6), Inches(1.2), Inches(5), Inches(2.5),
                "71%", font_size=84, color=RED, bold=True)

    add_textbox(slide, Inches(0.6), Inches(3.5), Inches(5), Inches(0.8),
                "of dementia patients have mixed\npathology at autopsy",
                font_size=18, color=TEXT_ON_LIGHT, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(4.4), Inches(3.5), color=RED)

    add_multiline_textbox(slide, Inches(0.6), Inches(4.7), Inches(5.5), Inches(2),
                          ["Clinical labels are inherently noisy. \"AD\" diagnoses are "
                           "confirmed at autopsy only 60-80% of the time.",
                           "",
                           "Models trained on uncertain labels develop uncertain "
                           "decision boundaries — garbage in, garbage out."],
                          font_size=13, color=TEXT_ON_LIGHT)

    # Right side context
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                "The Label Problem", font_size=22, color=RED, bold=True)

    add_accent_line(slide, Inches(7), Inches(2.1), Inches(2), color=RED)

    issues = [
        "Clinical diagnosis accuracy: 60-80% vs autopsy gold standard",
        "MCI is a syndrome, not a disease — heterogeneous etiology",
        "Mixed pathology (AD + vascular + Lewy body) is the norm, not the exception",
        "Longitudinal label changes: MCI → Normal reversion rate is 15-20%",
    ]

    for i, issue in enumerate(issues):
        y = Inches(2.5) + i * Inches(1.1)
        add_textbox(slide, Inches(7), y, Inches(5.5), Inches(0.9),
                    issue, font_size=13, color=TEXT_ON_LIGHT)
        if i < len(issues) - 1:
            add_shape_fill(slide, Inches(7), y + Inches(0.8), Inches(5), Pt(0.5), CITE_GRAY)

    add_citation_strip(slide,
        "Beach et al., 2012, J Neuropath Exp Neurol  •  Defined et al., 2016, Lancet Neurology  •  Defined et al., 2019, Brain")
    print("  Slide 27: Label Uncertainty")

    # ──────────────────────────────────────────────
    # SLIDE 28: "Recent Advances" (Split Compare, Light)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BG)

    add_section_label(slide, Inches(0.6), Inches(0.4), "MOVING FORWARD", color=GREEN)

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
                "Recent Advances", font_size=30, color=GREEN, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(2.5), color=GREEN)

    # Left: Multimodal fusion
    add_textbox(slide, Inches(0.6), Inches(2.2), Inches(5.5), Inches(0.5),
                "Multimodal Fusion", font_size=20, color=GREEN, bold=True)

    add_multiline_textbox(slide, Inches(0.6), Inches(2.8), Inches(5.5), Inches(3.5),
                          ["Combining MRI + PET + genetics + clinical scores yields "
                           "more robust predictions than any single modality.",
                           "",
                           "Early fusion (concatenate inputs) vs late fusion (combine predictions) "
                           "vs attention-based fusion (learn optimal weighting).",
                           "",
                           "Challenge: handling missing modalities gracefully — not every "
                           "patient has PET scans or genetic testing."],
                          font_size=13, color=TEXT_ON_LIGHT)

    # Divider
    add_shape_fill(slide, Inches(6.4), Inches(2.2), Pt(1.5), Inches(4.3), CITE_GRAY)

    # Right: Augmentation + Transfer
    add_textbox(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(0.5),
                "Augmentation & Transfer Learning", font_size=20, color=GREEN, bold=True)

    safe_add_image(slide, "Graph.png",
                   Inches(6.8), Inches(2.8), width=Inches(5.5), height=Inches(3.5))

    add_textbox(slide, Inches(6.8), Inches(6.3), Inches(5.5), Inches(0.5),
                "Transfer learning taxonomy and augmentation strategies",
                font_size=11, color=CITE_GRAY, italic=True)

    add_citation_strip(slide,
        "Defined et al., 2022, NeuroImage  •  Defined et al., 2023, Med Image Anal  •  Defined et al., 2021, IEEE TMI")
    print("  Slide 28: Recent Advances")

    # ──────────────────────────────────────────────
    # SLIDE 29: "Future Directions" (Dark Canvas)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_section_label(slide, Inches(0.6), Inches(0.5), "FUTURE", color=GREEN)

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
                "Future Directions", font_size=30, color=WHITE, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.8), Inches(2.5), color=GREEN)

    pillars = [
        ("Standardized\nBenchmarks", GREEN,
         "Unified preprocessing pipelines, common evaluation metrics, "
         "mandatory external validation. End the reproducibility crisis "
         "by making comparison possible."),
        ("Clinical\nInterpretability", BLUE,
         "Move beyond saliency maps to causal explanations. Clinicians need "
         "to understand not just 'where' but 'why'. Integrate explainability "
         "into model training, not post-hoc."),
        ("Multi-Stage\nProgression", COPPER,
         "Shift from binary AD/HC classification to predicting individual "
         "trajectories. Model the continuum from preclinical to severe. "
         "Personalized risk scoring over time."),
    ]

    for i, (title, color, desc) in enumerate(pillars):
        x = Inches(0.5) + i * Inches(4.2)
        box = add_shape_fill(slide, x, Inches(2.4), Inches(3.9), Inches(3.8),
                             RGBColor(0x16, 0x1B, 0x22))
        box.line.color.rgb = color
        box.line.width = Pt(2)

        add_textbox(slide, x + Inches(0.2), Inches(2.6), Inches(3.5), Inches(0.8),
                    title, font_size=20, color=color, bold=True)
        add_accent_line(slide, x + Inches(0.2), Inches(3.5), Inches(1.5), color=color)
        add_textbox(slide, x + Inches(0.2), Inches(3.7), Inches(3.5), Inches(2.2),
                    desc, font_size=12, color=TEXT_ON_DARK)

    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
                "The field must mature from 'can we classify?' to 'can we trust, deploy, and benefit patients?'",
                font_size=14, color=GREEN, italic=True,
                alignment=PP_ALIGN.CENTER)

    add_citation_strip(slide,
        "Defined et al., 2023, Nature Reviews Neuroscience  •  Defined et al., 2024, Lancet Digital Health",
        dark_bg=True)
    print("  Slide 29: Future Directions")

    # ──────────────────────────────────────────────
    # SLIDE 30: "Methodological Triad" (Statement, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                "The Methodological Triad", font_size=36, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(2.6), Inches(2.3), color=COPPER)

    triad = [
        ("1", "Subject-Level Splitting", "Never allow data from the same patient in both train and test sets. "
         "This single rule eliminates the most common source of inflated results."),
        ("2", "External Validation", "Test on independent datasets from different sites, scanners, and demographics. "
         "Internal cross-validation is necessary but insufficient."),
        ("3", "Confounder Control", "Account for age, sex, education, scanner effects, and head size. "
         "Without confounder analysis, classifiers may learn demographics rather than disease."),
    ]

    for i, (num, title, desc) in enumerate(triad):
        y = Inches(3.0) + i * Inches(1.4)
        add_textbox(slide, Inches(1.5), y, Inches(0.8), Inches(0.5),
                    num, font_size=32, color=COPPER, bold=True)
        add_textbox(slide, Inches(2.5), y, Inches(4), Inches(0.4),
                    title, font_size=18, color=WHITE, bold=True)
        add_textbox(slide, Inches(2.5), y + Inches(0.45), Inches(9), Inches(0.8),
                    desc, font_size=13, color=TEXT_ON_DARK)

    add_citation_strip(slide,
        "Wen et al., 2020, Med Image Anal  •  Defined et al., 2023, NeuroImage  •  Defined et al., 2022, Nature Methods",
        dark_bg=True)
    print("  Slide 30: Methodological Triad")

    # ──────────────────────────────────────────────
    # SLIDE 31: Conclusion (Statement, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(2),
                "Models are powerful.\nData is insufficient.\nTrust is unearned.",
                font_size=38, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(4.2), Inches(2.3), color=COPPER)

    add_textbox(slide, Inches(2), Inches(4.7), Inches(9.3), Inches(2),
                "The path from benchmark accuracy to clinical deployment requires not just "
                "better architectures, but better data practices, honest evaluation, and "
                "genuine collaboration between computer scientists and clinicians.",
                font_size=16, color=TEXT_ON_DARK, italic=True,
                alignment=PP_ALIGN.CENTER)

    add_citation_strip(slide, "Conclusion  •  AI Alzheimer and Dementia Classification", dark_bg=True)
    print("  Slide 31: Conclusion")

    # ──────────────────────────────────────────────
    # SLIDE 32: Thank You / Q&A (Custom, Dark)
    # ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    # Logo
    safe_add_image(slide, "logo_en.png", Inches(5.5), Inches(0.5),
                   height=Inches(1.2))

    add_accent_line(slide, Inches(5.5), Inches(2.0), Inches(2.3))

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                "Thank You", font_size=48, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
                "Questions & Discussion", font_size=22, color=COPPER, italic=True,
                alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(4.8), Inches(2.3))

    add_multiline_textbox(slide, Inches(3), Inches(5.2), Inches(7.3), Inches(1.5),
                          ["Gavriilidis Paraskevas",
                           "Department of Informatics and Computer Engineering",
                           "University of West Attica, 2025"],
                          font_size=14, color=CITE_GRAY,
                          alignment=PP_ALIGN.CENTER)

    add_citation_strip(slide,
        "AI Alzheimer and Dementia Classification  •  Thesis Defense  •  University of West Attica",
        dark_bg=True)
    print("  Slide 32: Thank You / Q&A")

    # ── Save ──
    prs.save(OUTPUT)
    print(f"\n✓ Presentation saved to: {OUTPUT}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
