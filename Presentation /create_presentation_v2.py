#!/usr/bin/env python3
"""
AI & Alzheimer's Disease Thesis Presentation v2
Condensed 24-slide presentation with experimental results from alzTheBatch.

Design system: Editorial magazine style, Spectral font, dark/light themes.
Slide dimensions: 13.3 x 7.5 inches
"""

import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# ─────────────────────────────────────────────────
# PATHS
# ─────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(os.path.dirname(BASE), "images")
ALZ_DIR = "/home/paris/Code/AlzTheBatch"
CHART_DIR = os.path.join(BASE, "generated_charts")
OUTPUT = os.path.join(BASE, "AI_Alzheimer_Thesis_Presentation_v2.pptx")
os.makedirs(CHART_DIR, exist_ok=True)

# ─────────────────────────────────────────────────
# DESIGN SYSTEM
# ─────────────────────────────────────────────────
# Colors
DARK_BG     = RGBColor(0x0D, 0x11, 0x17)
LIGHT_BG    = RGBColor(0xF7, 0xF5, 0xF0)
COPPER      = RGBColor(0xC1, 0x7F, 0x3A)
BLUE        = RGBColor(0x3B, 0x82, 0xB6)
RED         = RGBColor(0xDC, 0x4A, 0x4A)
GREEN       = RGBColor(0x5B, 0x8C, 0x6B)
TEXT_DARK   = RGBColor(0xE8, 0xE4, 0xDE)  # light text on dark bg
TEXT_LIGHT  = RGBColor(0x2D, 0x2D, 0x2D)  # dark text on light bg
CITATION_C  = RGBColor(0x8B, 0x86, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_ACCENT = RGBColor(0x1A, 0x1F, 0x2B)

# Hex for matplotlib (no # prefix for pptx, with # for matplotlib)
MPL_DARK_BG  = '#0D1117'
MPL_LIGHT_BG = '#F7F5F0'
MPL_COPPER   = '#C17F3A'
MPL_BLUE     = '#3B82B6'
MPL_RED      = '#DC4A4A'
MPL_GREEN    = '#5B8C6B'
MPL_TEXT_D   = '#E8E4DE'
MPL_TEXT_L   = '#2D2D2D'
MPL_CITATION = '#8B8680'

FONT = 'Spectral'
SLIDE_W = Inches(13.3)
SLIDE_H = Inches(7.5)

# Check if Spectral font is available, fall back to Georgia
_available_fonts = {f.name for f in fm.fontManager.ttflist}
MPL_FONT = 'Spectral' if 'Spectral' in _available_fonts else 'Georgia'


# ─────────────────────────────────────────────────
# CHART GENERATION
# ─────────────────────────────────────────────────

def setup_mpl():
    plt.rcParams.update({
        'font.family': 'serif',
        'font.serif': [MPL_FONT, 'Georgia', 'Times New Roman'],
        'font.size': 14,
        'axes.labelsize': 16,
        'axes.titlesize': 18,
        'xtick.labelsize': 13,
        'ytick.labelsize': 13,
        'figure.facecolor': 'none',
        'axes.facecolor': 'none',
        'savefig.facecolor': 'none',
        'savefig.transparent': True,
    })

def generate_prevalence_chart():
    """Slide 2: AD prevalence projection 2025-2060."""
    setup_mpl()
    years = [2025, 2030, 2035, 2040, 2045, 2050, 2055, 2060]
    millions = [7.2, 8.4, 9.6, 10.8, 11.7, 12.5, 13.2, 13.8]

    fig, ax = plt.subplots(figsize=(8, 4.5))
    ax.fill_between(years, millions, alpha=0.15, color=MPL_COPPER)
    ax.plot(years, millions, color=MPL_COPPER, linewidth=3, marker='o',
            markersize=8, markerfacecolor=MPL_COPPER, markeredgecolor='white', markeredgewidth=2)

    ax.set_xlabel('Year', color=MPL_TEXT_L, fontweight='bold')
    ax.set_ylabel('Millions of Americans', color=MPL_TEXT_L, fontweight='bold')
    ax.set_ylim(5, 15)
    ax.set_xlim(2023, 2062)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(MPL_CITATION)
    ax.spines['bottom'].set_color(MPL_CITATION)
    ax.tick_params(colors=MPL_TEXT_L)
    ax.yaxis.set_major_formatter(lambda x, _: f'{x:.0f}M')

    for yr, val in zip(years[::2], millions[::2]):
        ax.annotate(f'{val}M', (yr, val), textcoords="offset points",
                    xytext=(0, 14), ha='center', fontsize=12,
                    fontweight='bold', color=MPL_COPPER)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'prevalence_projection.png')
    plt.savefig(path, dpi=300, transparent=True, bbox_inches='tight')
    plt.close()
    return path

def generate_classical_ml_chart():
    """Slide 13: Classical ML comparison -- AD-vs-HC vs MCI."""
    setup_mpl()
    categories = ['AD vs HC\n(Binary)', 'MCI Detection\n(Multi-class)']
    svm_vals = [94.5, 68.0]
    rf_vals = [91.0, 62.0]

    fig, ax = plt.subplots(figsize=(7, 4.5))
    x = np.arange(len(categories))
    width = 0.3

    bars1 = ax.bar(x - width/2, svm_vals, width, label='SVM', color=MPL_BLUE,
                   edgecolor='white', linewidth=1.5, zorder=3)
    bars2 = ax.bar(x + width/2, rf_vals, width, label='Random Forest', color=MPL_COPPER,
                   edgecolor='white', linewidth=1.5, zorder=3)

    ax.set_ylabel('Accuracy (%)', color=MPL_TEXT_L, fontweight='bold')
    ax.set_ylim(0, 105)
    ax.set_xticks(x)
    ax.set_xticklabels(categories, color=MPL_TEXT_L)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(MPL_CITATION)
    ax.spines['bottom'].set_color(MPL_CITATION)
    ax.tick_params(colors=MPL_TEXT_L)
    ax.legend(frameon=False, fontsize=12)

    for bars in [bars1, bars2]:
        for bar in bars:
            h = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., h + 1.5,
                    f'{h:.1f}%', ha='center', va='bottom',
                    fontweight='bold', fontsize=13, color=MPL_TEXT_L)

    # Gap annotation
    ax.annotate('', xy=(1.15, 68), xytext=(1.15, 94.5),
                arrowprops=dict(arrowstyle='<->', color=MPL_RED, lw=2))
    ax.text(1.35, 81, '~27%\ngap', ha='center', va='center',
            fontsize=12, fontweight='bold', color=MPL_RED)

    ax.axhline(y=70, color=MPL_CITATION, linestyle='--', alpha=0.4, linewidth=1)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'classical_ml_comparison.png')
    plt.savefig(path, dpi=300, transparent=True, bbox_inches='tight')
    plt.close()
    return path

def generate_cnn_vs_swin_chart():
    """Slide 18: CNN vs Swin head-to-head comparison."""
    setup_mpl()
    metrics = ['Accuracy', 'Balanced\nAccuracy', 'F1 Macro', 'MCC']
    cnn_vals = [59.3, 63.5, 51.9, 35.6]
    swin_vals = [87.6, 90.4, 89.2, 80.0]

    fig, ax = plt.subplots(figsize=(8, 4.5))
    x = np.arange(len(metrics))
    width = 0.32

    bars1 = ax.bar(x - width/2, cnn_vals, width, label='CNN (Combined)',
                   color=MPL_RED, edgecolor='white', linewidth=1.5, zorder=3)
    bars2 = ax.bar(x + width/2, swin_vals, width, label='Swin Transformer',
                   color=MPL_GREEN, edgecolor='white', linewidth=1.5, zorder=3)

    ax.set_ylabel('Score (%)', color=MPL_TEXT_L, fontweight='bold')
    ax.set_ylim(0, 105)
    ax.set_xticks(x)
    ax.set_xticklabels(metrics, color=MPL_TEXT_L)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(MPL_CITATION)
    ax.spines['bottom'].set_color(MPL_CITATION)
    ax.tick_params(colors=MPL_TEXT_L)
    ax.legend(frameon=False, fontsize=13, loc='upper left')

    for bars in [bars1, bars2]:
        for bar in bars:
            h = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., h + 1.5,
                    f'{h:.1f}%', ha='center', va='bottom',
                    fontweight='bold', fontsize=12, color=MPL_TEXT_L)

    # Improvement arrows
    for i in range(len(metrics)):
        diff = swin_vals[i] - cnn_vals[i]
        mid_x = x[i]
        ax.annotate(f'+{diff:.0f}pp', xy=(mid_x, max(swin_vals[i], cnn_vals[i]) + 8),
                    ha='center', fontsize=10, fontweight='bold', color=MPL_GREEN)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'cnn_vs_swin_comparison.png')
    plt.savefig(path, dpi=300, transparent=True, bbox_inches='tight')
    plt.close()
    return path

def generate_leakage_chart():
    """Slide 19: Data leakage impact chart."""
    setup_mpl()
    categories = ['With Data\nLeakage', 'Proper\nMethodology']
    values = [95.0, 67.0]
    colors = [MPL_RED, MPL_GREEN]

    fig, ax = plt.subplots(figsize=(6, 4.5))
    bars = ax.bar(categories, values, width=0.5, color=colors,
                  edgecolor='white', linewidth=2, zorder=3)

    ax.set_ylabel('Reported Accuracy (%)', color=MPL_TEXT_L, fontweight='bold')
    ax.set_ylim(0, 110)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(MPL_CITATION)
    ax.spines['bottom'].set_color(MPL_CITATION)
    ax.tick_params(colors=MPL_TEXT_L)

    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2., val + 2,
                f'{val:.0f}%', ha='center', va='bottom',
                fontweight='bold', fontsize=18, color=MPL_TEXT_L)

    # Drop annotation
    ax.annotate('', xy=(1, 67), xytext=(0, 95),
                arrowprops=dict(arrowstyle='->', color=MPL_TEXT_L, lw=2.5,
                                connectionstyle='arc3,rad=-0.3'))
    ax.text(0.5, 78, '-28pp', ha='center', fontsize=16,
            fontweight='bold', color=MPL_RED)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, 'data_leakage_impact.png')
    plt.savefig(path, dpi=300, transparent=True, bbox_inches='tight')
    plt.close()
    return path


# ─────────────────────────────────────────────────
# PRESENTATION HELPERS
# ─────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=16,
                 color=TEXT_LIGHT, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a text box with multiple formatted runs in a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    for i, run_data in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = run_data.get('text', '')
        run.font.name = run_data.get('font', FONT)
        run.font.size = Pt(run_data.get('size', 16))
        run.font.color.rgb = run_data.get('color', TEXT_LIGHT)
        run.font.bold = run_data.get('bold', False)
        run.font.italic = run_data.get('italic', False)
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=14,
                       color=TEXT_LIGHT, font_name=FONT, bold=False,
                       alignment=PP_ALIGN.LEFT, line_spacing=None, anchor=MSO_ANCHOR.TOP):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, dict):
            p.text = line.get('text', '')
            p.font.size = Pt(line.get('size', font_size))
            p.font.color.rgb = line.get('color', color)
            p.font.bold = line.get('bold', bold)
            p.font.italic = line.get('italic', False)
            p.font.name = line.get('font', font_name)
            p.alignment = line.get('alignment', alignment)
            if line.get('spacing'):
                p.space_before = Pt(line['spacing'])
        else:
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = font_name
            p.alignment = alignment
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox

def add_citation(slide, text, bg_dark=True):
    """Add citation strip at bottom of slide."""
    color = CITATION_C
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                 text, font_size=9, color=color, italic=True)

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if file exists, skip gracefully if not."""
    if os.path.exists(path):
        kwargs = {'left': left, 'top': top}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    else:
        print(f"  WARNING: Image not found: {path}")
        return False

def add_accent_line(slide, left, top, width, color=COPPER):
    """Add a thin decorative accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ─────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Logo
    logo = os.path.join(IMG_DIR, "logo_en.png")
    add_image_safe(slide, logo, Inches(0.6), Inches(0.4), width=Inches(2.8))

    # Accent line
    add_accent_line(slide, Inches(0.6), Inches(3.2), Inches(3))

    # Title
    add_text_box(slide, Inches(0.6), Inches(3.5), Inches(11), Inches(1.2),
                 "Artificial Intelligence in Alzheimer's Disease Diagnosis",
                 font_size=34, color=TEXT_DARK, bold=True)

    # Subtitle
    add_text_box(slide, Inches(0.6), Inches(4.8), Inches(11), Inches(0.8),
                 "From Neuroimaging Pipelines to Deep Learning Classification",
                 font_size=20, color=COPPER, italic=True)

    # Author & University
    add_multiline_text(slide, Inches(0.6), Inches(5.8), Inches(6), Inches(1.2), [
        {'text': 'Paris Karageorgakis', 'size': 16, 'color': TEXT_DARK, 'bold': True},
        {'text': 'University of Piraeus, Department of Informatics', 'size': 13, 'color': CITATION_C, 'spacing': 6},
        {'text': 'Thesis Advisor: Prof. Christos Douligeris', 'size': 13, 'color': CITATION_C, 'spacing': 4},
    ])


def slide_02_epidemic(prs, chart_path):
    """The Silent Epidemic -- hero stat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    # Section tag
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.4),
                 "THE PROBLEM", font_size=11, color=COPPER, bold=True)

    # Hero number
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(5), Inches(1.5),
                 "7.2 Million", font_size=72, color=DARK_BG, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.4), Inches(5), Inches(0.5),
                 "Americans living with Alzheimer's disease", font_size=18,
                 color=TEXT_LIGHT, bold=False)

    # Supporting stats
    stats = [
        "Every 65 seconds, someone develops AD",
        "$360 billion annual healthcare cost",
        "6th leading cause of death in the US",
        "Projected to reach 13.8M by 2060",
    ]
    add_multiline_text(slide, Inches(0.8), Inches(3.2), Inches(5), Inches(2.5), [
        {'text': f'\u2022  {s}', 'size': 14, 'color': TEXT_LIGHT, 'spacing': 8}
        for s in stats
    ], line_spacing=22)

    # Chart on right
    add_image_safe(slide, chart_path, Inches(6.5), Inches(1.0), width=Inches(6.2))

    add_citation(slide, "Alzheimer's Association, 2024 Facts and Figures  |  Rajan et al., 2021", bg_dark=False)


def slide_03_window_atn(prs):
    """The 20-Year Window & AT(N) -- split compare."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Left half - dark
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "THE 20-YEAR WINDOW", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(5.8), Inches(1.5),
                 "Pathological changes begin 15\u201320 years before the first clinical symptoms.",
                 font_size=26, color=TEXT_DARK, bold=True)

    add_multiline_text(slide, Inches(0.6), Inches(2.8), Inches(5.8), Inches(3.5), [
        {'text': 'By the time of diagnosis, neuronal loss is already irreversible.',
         'size': 15, 'color': TEXT_DARK},
        {'text': '', 'size': 8, 'color': TEXT_DARK},
        {'text': 'The diagnostic challenge:', 'size': 15, 'color': COPPER, 'bold': True, 'spacing': 12},
        {'text': '\u2022  MCI patients convert to AD at 10\u201315% per year',
         'size': 14, 'color': TEXT_DARK, 'spacing': 8},
        {'text': '\u2022  Clinical diagnosis accuracy: ~77% (confirmed at autopsy)',
         'size': 14, 'color': TEXT_DARK, 'spacing': 6},
        {'text': '\u2022  Early intervention could slow progression by 30%',
         'size': 14, 'color': TEXT_DARK, 'spacing': 6},
    ], line_spacing=20)

    # Vertical divider
    add_shape_rect(slide, Inches(6.6), Inches(0.5), Pt(2), Inches(6.2), COPPER)

    # Right half - AT(N) Framework
    add_text_box(slide, Inches(7.0), Inches(0.4), Inches(6), Inches(0.4),
                 "AT(N) FRAMEWORK", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(7.0), Inches(1.0), Inches(5.5), Inches(0.8),
                 "Biological Classification of AD", font_size=22, color=TEXT_DARK, bold=True)

    # AT(N) boxes
    atn = [
        ("A", "Amyloid", "A\u03b2 plaques accumulate\n15\u201320 years before symptoms", BLUE),
        ("T", "Tau", "Neurofibrillary tangles spread\nalong predictable pathways", COPPER),
        ("(N)", "Neurodegeneration", "Synaptic loss, brain atrophy\nmeasurable on MRI", RED),
    ]
    y_pos = Inches(2.0)
    for letter, label, desc, color in atn:
        # Letter box
        box = add_rounded_rect(slide, Inches(7.0), y_pos, Inches(0.9), Inches(1.2), color)
        add_text_box(slide, Inches(7.0), y_pos + Inches(0.15), Inches(0.9), Inches(0.9),
                     letter, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Label and description
        add_text_box(slide, Inches(8.1), y_pos + Inches(0.05), Inches(4.5), Inches(0.4),
                     label, font_size=16, color=TEXT_DARK, bold=True)
        add_text_box(slide, Inches(8.1), y_pos + Inches(0.45), Inches(4.5), Inches(0.7),
                     desc, font_size=12, color=CITATION_C)
        y_pos += Inches(1.5)

    add_citation(slide, "Jack et al., 2018  |  Sperling et al., 2011  |  NIA-AA Research Framework")


def slide_04_neuroimaging(prs):
    """Why Neuroimaging -- full-bleed image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Image (takes most of the slide)
    img = os.path.join(IMG_DIR, "nihms-137059-f0004.jpg")
    add_image_safe(slide, img, Inches(5.5), Inches(0.3), width=Inches(7.5))

    # Dark overlay panel on left
    add_shape_rect(slide, Inches(0), Inches(0), Inches(6), SLIDE_H, DARK_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(5), Inches(0.4),
                 "WHY NEUROIMAGING", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(5), Inches(1.2),
                 "The Window into\nBrain Pathology", font_size=32, color=TEXT_DARK, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(2.5), Inches(2))

    benefits = [
        ("Non-invasive", "No surgery, no lumbar puncture required"),
        ("Quantifiable", "Volumetric measurements enable objective tracking"),
        ("Reproducible", "Standardized protocols across clinical sites"),
        ("Accessible", "MRI available in most medical centers worldwide"),
    ]
    y = Inches(3.0)
    for title, desc in benefits:
        add_text_box(slide, Inches(0.6), y, Inches(5), Inches(0.35),
                     title, font_size=16, color=COPPER, bold=True)
        add_text_box(slide, Inches(0.6), y + Inches(0.35), Inches(5), Inches(0.4),
                     desc, font_size=13, color=TEXT_DARK)
        y += Inches(0.95)

    add_citation(slide, "Jack et al., 2010  |  Defined in MRI context by  Defined in MRI context")


def slide_05_datasets(prs):
    """The Benchmark Datasets -- 3-column split."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(5), Inches(0.4),
                 "BENCHMARK DATASETS", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(10), Inches(0.8),
                 "The Three Pillars of AD Neuroimaging Research",
                 font_size=28, color=TEXT_LIGHT, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.8), Inches(2))

    datasets = [
        {
            'name': 'ADNI',
            'full': 'Alzheimer\'s Disease\nNeuroimaging Initiative',
            'color': BLUE,
            'stats': ['2,400+ subjects', 'Longitudinal since 2004', 'MRI + PET + CSF + genetics',
                      'Gold standard for method validation'],
        },
        {
            'name': 'AIBL',
            'full': 'Australian Imaging,\nBiomarkers & Lifestyle',
            'color': COPPER,
            'stats': ['1,100+ subjects', 'Melbourne & Perth cohorts', 'Focus on lifestyle factors',
                      'Diverse population sampling'],
        },
        {
            'name': 'OASIS',
            'full': 'Open Access Series\nof Imaging Studies',
            'color': GREEN,
            'stats': ['2,000+ subjects', 'Cross-sectional + longitudinal', 'Freely available',
                      'Ages 18\u201396 years'],
        },
    ]

    x_positions = [Inches(0.6), Inches(4.7), Inches(8.8)]
    for i, ds in enumerate(datasets):
        x = x_positions[i]

        # Header box
        add_rounded_rect(slide, x, Inches(2.2), Inches(3.6), Inches(1.2), ds['color'])
        add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(3.2), Inches(0.5),
                     ds['name'], font_size=24, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(3.2), Inches(0.5),
                     ds['full'], font_size=11, color=RGBColor(0xFF, 0xFF, 0xFF))

        # Stats
        y = Inches(3.7)
        for stat in ds['stats']:
            add_text_box(slide, x + Inches(0.2), y, Inches(3.4), Inches(0.35),
                         f'\u2022  {stat}', font_size=13, color=TEXT_LIGHT)
            y += Inches(0.4)

    add_citation(slide, "Mueller et al., 2005 (ADNI)  |  Ellis et al., 2009 (AIBL)  |  Marcus et al., 2007 (OASIS)", bg_dark=False)


def slide_06_section_divider(prs):
    """Section divider: 'Seeing the Brain'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
                 "ACT II", font_size=11, color=COPPER, bold=True)

    add_accent_line(slide, Inches(0.8), Inches(2.8), Inches(4))

    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(1.5),
                 "Seeing the Brain",
                 font_size=52, color=TEXT_DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(10), Inches(0.8),
                 "From MRI physics to preprocessing pipelines to machine learning",
                 font_size=18, color=CITATION_C, italic=True)


def slide_07_mri_fundamentals(prs):
    """MRI Fundamentals."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Image on right
    img = os.path.join(IMG_DIR, "IntensityNormalization1.png")
    add_image_safe(slide, img, Inches(7.0), Inches(0.5), width=Inches(5.8))

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "IMAGING FUNDAMENTALS", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(6), Inches(1.0),
                 "MRI: The Structural Gold Standard",
                 font_size=28, color=TEXT_DARK, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(2.0), Inches(2))

    points = [
        "Nuclear magnetic resonance of hydrogen atoms",
        "T1-weighted: Gray matter dark, white matter bright",
        "T2-weighted: CSF bright, sensitive to edema",
        "Sub-millimeter resolution (0.5\u20131 mm voxels)",
        "Larmor equation: \u03c9 = \u03b3B\u2080 governs precession frequency",
        "No ionizing radiation \u2014 safe for longitudinal studies",
    ]
    y = Inches(2.4)
    for pt in points:
        add_text_box(slide, Inches(0.6), y, Inches(6.2), Inches(0.4),
                     f'\u2022  {pt}', font_size=14, color=TEXT_DARK)
        y += Inches(0.55)

    add_citation(slide, "Bitar et al., 2006  |  Symms et al., 2004  |  McRobbie et al., 2017")


def slide_08_beyond_mri(prs):
    """Beyond MRI: CT & PET (merged from original 9+10)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "IMAGING MODALITIES BEYOND MRI", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(10), Inches(0.7),
                 "CT & PET: Complementary Windows", font_size=28, color=TEXT_LIGHT, bold=True)

    # Left: CT
    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(0.5),
                 "Computed Tomography", font_size=18, color=BLUE, bold=True)

    ct_img = os.path.join(IMG_DIR, "pmp-32-1-1-f1.png")
    add_image_safe(slide, ct_img, Inches(0.6), Inches(2.4), width=Inches(3.0))

    ct_points = ["X-ray attenuation imaging", "Fast acquisition (~seconds)",
                 "Detects gross atrophy & calcifications", "Limited soft-tissue contrast vs MRI"]
    y = Inches(2.5)
    for pt in ct_points:
        add_text_box(slide, Inches(3.8), y, Inches(2.6), Inches(0.35),
                     f'\u2022  {pt}', font_size=12, color=TEXT_LIGHT)
        y += Inches(0.42)

    # Divider
    add_shape_rect(slide, Inches(6.5), Inches(1.8), Pt(2), Inches(4.8), COPPER)

    # Right: PET
    add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.5),
                 "PET Biomarkers", font_size=18, color=COPPER, bold=True)

    pet_img = os.path.join(IMG_DIR, "pmp-32-1-1-f4.png")
    add_image_safe(slide, pet_img, Inches(7.0), Inches(2.4), width=Inches(3.0))

    pet_points = [
        "FDG-PET: 90% sensitivity for AD",
        "Amyloid PET: Detects A\u03b2 plaques in vivo",
        "Tau PET: Maps tangle distribution",
        "Quantifies molecular pathology directly",
    ]
    y = Inches(2.5)
    for pt in pet_points:
        add_text_box(slide, Inches(10.2), y, Inches(2.6), Inches(0.35),
                     f'\u2022  {pt}', font_size=12, color=TEXT_LIGHT)
        y += Inches(0.42)

    add_citation(slide, "Marcus et al., 2014  |  Johnson et al., 2012  |  Minoshima et al., 1997  |  Clark et al., 2011", bg_dark=False)


def slide_09_preprocessing(prs):
    """The Preprocessing Pipeline -- dark canvas with pipeline flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "THE PREPROCESSING PIPELINE", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.7),
                 "5 Steps from Raw Scan to Analysis-Ready Data",
                 font_size=28, color=TEXT_DARK, bold=True)

    steps = [
        ("01", "Registration", "Align to standard space\n(MNI152 template)", BLUE),
        ("02", "Skull Stripping", "Remove non-brain tissue\n(BET, HD-BET, SynthStrip)", COPPER),
        ("03", "Intensity\nNormalization", "Standardize signal\n(Z-Score, White Stripe)", GREEN),
        ("04", "Denoising", "Remove acquisition noise\n(NLM, BM3D, Deep Learning)", RED),
        ("05", "Segmentation", "Classify tissue types\n(GM, WM, CSF)", BLUE),
    ]

    x = Inches(0.4)
    for num, title, desc, color in steps:
        # Box
        add_rounded_rect(slide, x, Inches(2.5), Inches(2.3), Inches(3.8), DARK_ACCENT, color)

        # Step number
        add_text_box(slide, x, Inches(2.6), Inches(2.3), Inches(0.7),
                     num, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.15), Inches(3.4), Inches(2.0), Inches(0.8),
                     title, font_size=15, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + Inches(0.15), Inches(4.3), Inches(2.0), Inches(1.2),
                     desc, font_size=11, color=CITATION_C, alignment=PP_ALIGN.CENTER)

        # Arrow (except last)
        if num != "05":
            add_text_box(slide, x + Inches(2.25), Inches(3.8), Inches(0.35), Inches(0.5),
                         '\u2192', font_size=24, color=COPPER, bold=True, alignment=PP_ALIGN.CENTER)

        x += Inches(2.55)

    add_citation(slide, "Ashburner, 2012  |  Smith, 2002  |  Manjón & Coupé, 2016")


def slide_10_signal_cleaning(prs):
    """Signal Cleaning (merged intensity norm + denoising)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "SIGNAL CLEANING", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.7),
                 "Intensity Normalization & Denoising", font_size=28, color=TEXT_LIGHT, bold=True)

    # Left: Intensity Normalization
    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(0.5),
                 "Intensity Normalization", font_size=18, color=BLUE, bold=True)

    img2 = os.path.join(IMG_DIR, "IntensityNormalization2.png")
    add_image_safe(slide, img2, Inches(0.6), Inches(2.4), width=Inches(5.8))

    norm_points = ["Z-Score: Mean-center, unit-variance per subject",
                   "White Stripe: Normalize to normal-appearing white matter",
                   "Essential for cross-site comparisons"]
    y = Inches(5.0)
    for pt in norm_points:
        add_text_box(slide, Inches(0.6), y, Inches(5.8), Inches(0.35),
                     f'\u2022  {pt}', font_size=12, color=TEXT_LIGHT)
        y += Inches(0.38)

    # Divider
    add_shape_rect(slide, Inches(6.6), Inches(1.8), Pt(2), Inches(4.8), COPPER)

    # Right: Denoising
    add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.5),
                 "Denoising", font_size=18, color=COPPER, bold=True)

    img3 = os.path.join(IMG_DIR, "IntensityNormalization3.png")
    add_image_safe(slide, img3, Inches(7.0), Inches(2.4), width=Inches(5.8))

    denoise_points = ["NLM: Non-Local Means (patch-based similarity)",
                      "BM3D: Block-Matching 3D (transform-domain filtering)",
                      "Deep Learning: CNN-based denoising autoencoders"]
    y = Inches(5.0)
    for pt in denoise_points:
        add_text_box(slide, Inches(7.0), y, Inches(5.8), Inches(0.35),
                     f'\u2022  {pt}', font_size=12, color=TEXT_LIGHT)
        y += Inches(0.38)

    add_citation(slide, "Shinohara et al., 2014  |  Buades et al., 2005  |  Dabov et al., 2007", bg_dark=False)


def slide_11_skull_stripping(prs):
    """Skull Stripping."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Background image
    img1 = os.path.join(IMG_DIR, "Skull Stripping image.png")
    add_image_safe(slide, img1, Inches(6.5), Inches(0.3), width=Inches(6.5))

    # Left panel overlay
    add_shape_rect(slide, Inches(0), Inches(0), Inches(7), SLIDE_H, DARK_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "SKULL STRIPPING", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(6), Inches(1.0),
                 "Removing Non-Brain Tissue", font_size=28, color=TEXT_DARK, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(2.0), Inches(2))

    methods = [
        ("BET", "Brain Extraction Tool \u2014 surface deformation model"),
        ("HD-BET", "Deep learning \u2014 robust across scanners, 95%+ Dice"),
        ("SynthStrip", "Synthesis-based \u2014 contrast-agnostic extraction"),
    ]
    y = Inches(2.5)
    for name, desc in methods:
        add_text_box(slide, Inches(0.6), y, Inches(2), Inches(0.4),
                     name, font_size=16, color=COPPER, bold=True)
        add_text_box(slide, Inches(2.8), y, Inches(3.6), Inches(0.4),
                     desc, font_size=13, color=TEXT_DARK)
        y += Inches(0.55)

    # Warning box
    add_rounded_rect(slide, Inches(0.6), Inches(4.5), Inches(5.8), Inches(1.5),
                     DARK_ACCENT, RED)
    add_text_box(slide, Inches(0.9), Inches(4.6), Inches(5.2), Inches(0.4),
                 '\u26a0  Shortcut Learning Risk', font_size=14, color=RED, bold=True)
    add_text_box(slide, Inches(0.9), Inches(5.0), Inches(5.2), Inches(0.8),
                 'Models can exploit skull artifacts as class features rather than learning true brain pathology patterns.',
                 font_size=12, color=TEXT_DARK)

    # Techniques image
    img2 = os.path.join(IMG_DIR, "Skull Stripping Techniques.png")
    add_image_safe(slide, img2, Inches(0.6), Inches(6.1), width=Inches(5.8))

    add_citation(slide, "Smith, 2002  |  Isensee et al., 2019  |  Hoopes et al., 2022")


def slide_12_vbm(prs):
    """Voxel-Based Morphometry."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Full-bleed image
    img = os.path.join(IMG_DIR, "nihms154848f1.jpg")
    add_image_safe(slide, img, Inches(5.5), Inches(0), width=Inches(7.8))

    # Left panel
    add_shape_rect(slide, Inches(0), Inches(0), Inches(6.2), SLIDE_H, DARK_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(5), Inches(0.4),
                 "VOXEL-BASED MORPHOMETRY", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(5.4), Inches(1.0),
                 "Whole-Brain Analysis\nof Structural Changes",
                 font_size=28, color=TEXT_DARK, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(2.3), Inches(2))

    points = [
        "Voxel-by-voxel statistical comparison of gray matter density",
        "Hippocampal atrophy: hallmark signature of early AD",
        "AUC > 0.90 for AD vs healthy controls",
        "Reveals distributed patterns invisible to visual inspection",
        "Can track atrophy progression over time",
    ]
    y = Inches(2.8)
    for pt in points:
        add_text_box(slide, Inches(0.6), y, Inches(5.4), Inches(0.5),
                     f'\u2022  {pt}', font_size=14, color=TEXT_DARK)
        y += Inches(0.6)

    add_citation(slide, "Ashburner & Friston, 2000  |  Karas et al., 2004  |  Whitwell, 2009")


def slide_13_classical_ml(prs, chart_path):
    """Classical ML: The Ceiling -- hero stat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "CLASSICAL MACHINE LEARNING", font_size=11, color=COPPER, bold=True)

    # Hero stat
    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(5), Inches(1.5),
                 "94.5%", font_size=72, color=BLUE, bold=True)

    add_text_box(slide, Inches(0.6), Inches(2.3), Inches(5), Inches(0.5),
                 "SVM accuracy for AD vs. healthy controls", font_size=18, color=TEXT_LIGHT)

    add_text_box(slide, Inches(0.6), Inches(3.0), Inches(5), Inches(0.5),
                 "But MCI detection plummets to ~68%", font_size=16, color=RED, bold=True)

    points = [
        "SVM, Random Forest, Logistic Regression dominate pre-2016",
        "Handcrafted features: volume, cortical thickness, texture",
        "Binary AD/HC: near-clinical performance",
        "Multi-class with MCI: dramatic accuracy drop",
        "Feature engineering = manual, limited, domain-dependent",
    ]
    y = Inches(3.7)
    for pt in points:
        add_text_box(slide, Inches(0.6), y, Inches(5.5), Inches(0.4),
                     f'\u2022  {pt}', font_size=13, color=TEXT_LIGHT)
        y += Inches(0.45)

    # Chart on right
    add_image_safe(slide, chart_path, Inches(6.5), Inches(0.8), width=Inches(6.3))

    add_citation(slide, "Klöppel et al., 2008  |  Cuingnet et al., 2011  |  Salvatore et al., 2015", bg_dark=False)


def slide_14_cnns_to_transformers(prs):
    """From CNNs to Transformers (merged)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "DEEP LEARNING ARCHITECTURES", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.7),
                 "From CNNs to Transformers", font_size=28, color=TEXT_LIGHT, bold=True)

    # Left column: CNNs
    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(0.5),
                 "Convolutional Neural Networks", font_size=18, color=BLUE, bold=True)

    cnn_items = [
        ("2D CNNs", "Single-slice processing (ResNet-18, VGG-16)"),
        ("3D CNNs", "Volumetric convolutions capture spatial context"),
        ("DenseNet", "Dense connections, feature reuse, fewer parameters"),
        ("ResNet", "Skip connections enable very deep networks (152+ layers)"),
        ("Hybrid CNN+SVM", "CNN features fed to SVM classifier \u2014 82\u201390% on AD"),
    ]
    y = Inches(2.4)
    for title, desc in cnn_items:
        add_rich_text(slide, Inches(0.6), y, Inches(5.8), Inches(0.5), [
            {'text': f'{title}:  ', 'size': 13, 'color': TEXT_LIGHT, 'bold': True},
            {'text': desc, 'size': 13, 'color': TEXT_LIGHT},
        ])
        y += Inches(0.5)

    # Divider
    add_shape_rect(slide, Inches(6.5), Inches(1.8), Pt(2), Inches(4.8), COPPER)

    # Right column: Transformers
    add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.5),
                 "Vision Transformers", font_size=18, color=COPPER, bold=True)

    tx_items = [
        ("ViT", "Patches as tokens, global attention, data-hungry"),
        ("Swin Transformer", "Shifted windows \u2014 local + global attention, efficient"),
        ("DeiT", "Data-efficient training with knowledge distillation"),
        ("Hybrid ViT+CNN", "CNN backbone + transformer head for small datasets"),
        ("Key advantage", "Pre-training on ImageNet transfers to medical imaging"),
    ]
    y = Inches(2.4)
    for title, desc in tx_items:
        add_rich_text(slide, Inches(7.0), y, Inches(5.8), Inches(0.5), [
            {'text': f'{title}:  ', 'size': 13, 'color': TEXT_LIGHT, 'bold': True},
            {'text': desc, 'size': 13, 'color': TEXT_LIGHT},
        ])
        y += Inches(0.5)

    # Bottom insight box
    add_rounded_rect(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.2),
                     RGBColor(0xED, 0xEB, 0xE5))
    add_rich_text(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.8), [
        {'text': 'Key Insight: ', 'size': 14, 'color': COPPER, 'bold': True},
        {'text': 'Transformers achieve state-of-the-art performance on medical imaging benchmarks, '
                 'but require careful handling of small datasets and class imbalance '
                 '\u2014 motivating our experimental validation in the next section.',
         'size': 14, 'color': TEXT_LIGHT},
    ])

    add_citation(slide, "Dosovitskiy et al., 2021  |  Liu et al., 2021  |  He et al., 2016  |  Huang et al., 2017", bg_dark=False)


def slide_15_explainability(prs):
    """Explainability & Grad-CAM."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "EXPLAINABILITY (XAI)", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(6), Inches(0.7),
                 "Opening the Black Box", font_size=28, color=TEXT_DARK, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(2))

    pillars = [
        ("Transparency", "How does the model make decisions?"),
        ("Interpretability", "Can clinicians understand the reasoning?"),
        ("Trustworthiness", "Are predictions reliable for clinical use?"),
    ]
    y = Inches(2.1)
    for title, desc in pillars:
        add_rich_text(slide, Inches(0.6), y, Inches(5.8), Inches(0.4), [
            {'text': f'{title}: ', 'size': 14, 'color': COPPER, 'bold': True},
            {'text': desc, 'size': 14, 'color': TEXT_DARK},
        ])
        y += Inches(0.45)

    methods = ["Grad-CAM: Gradient-weighted class activation maps",
               "LIME: Local interpretable model-agnostic explanations",
               "SHAP: SHapley Additive exPlanations (game-theoretic)"]
    y = Inches(3.6)
    for m in methods:
        add_text_box(slide, Inches(0.6), y, Inches(5.8), Inches(0.35),
                     f'\u2022  {m}', font_size=13, color=TEXT_DARK)
        y += Inches(0.42)

    # Images
    img1 = os.path.join(IMG_DIR, "Grad-CAMVBM.png")
    img2 = os.path.join(IMG_DIR, "limitations_gradcam.png")
    add_image_safe(slide, img1, Inches(6.8), Inches(0.4), width=Inches(6.0))
    add_image_safe(slide, img2, Inches(6.8), Inches(3.8), width=Inches(6.0))

    add_citation(slide, "Selvaraju et al., 2017  |  Ribeiro et al., 2016  |  Lundberg & Lee, 2017")


def slide_16_cnn_experiments(prs):
    """NEW: Our CNN Experiments."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "OUR EXPERIMENTS", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(6), Inches(0.7),
                 "CNN: 5 Approaches to Class Imbalance",
                 font_size=26, color=TEXT_DARK, bold=True)

    # Key finding box
    add_rounded_rect(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(1.3),
                     DARK_ACCENT, GREEN)
    add_multiline_text(slide, Inches(0.9), Inches(1.9), Inches(5.2), Inches(1.1), [
        {'text': 'Combined Strategy wins:', 'size': 14, 'color': GREEN, 'bold': True},
        {'text': 'Class weights + balanced sampling + MONAI augmentation', 'size': 12, 'color': TEXT_DARK, 'spacing': 4},
        {'text': '63.5% balanced accuracy  |  100% Moderate recall  |  F1: 0.52', 'size': 12, 'color': COPPER, 'bold': True, 'spacing': 4},
    ])

    # Dataset info
    add_multiline_text(slide, Inches(0.6), Inches(3.4), Inches(5.8), Inches(1.5), [
        {'text': 'Dataset: 11,519 MRI scans (Falah/Alzheimer_MRI)', 'size': 12, 'color': CITATION_C},
        {'text': '4 classes: Non Dem. (3200) | Very Mild (3008) | Mild (2739) | Moderate (2572)', 'size': 11, 'color': CITATION_C, 'spacing': 4},
        {'text': '', 'size': 6, 'color': TEXT_DARK},
        {'text': 'Baseline CNN: 0% recall on Mild & Moderate classes', 'size': 13, 'color': RED, 'bold': True, 'spacing': 6},
        {'text': 'Standard training fails for imbalanced medical data.', 'size': 13, 'color': TEXT_DARK, 'spacing': 4},
    ])

    # Training curves image
    img1 = os.path.join(ALZ_DIR, "training_curves_comparison.png")
    add_image_safe(slide, img1, Inches(6.8), Inches(0.3), width=Inches(6.0))

    # Confusion matrices
    img2 = os.path.join(ALZ_DIR, "confusion_matrices_comparison.png")
    add_image_safe(slide, img2, Inches(6.8), Inches(3.8), width=Inches(6.0))

    add_citation(slide, "Author's experimental results, alzTheBatch repository  |  HuggingFace: Falah/Alzheimer_MRI  |  MONAI Framework")


def slide_17_swin_results(prs):
    """NEW: Swin Transformer Results -- hero stat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "SWIN TRANSFORMER RESULTS", font_size=11, color=COPPER, bold=True)

    # Hero stat
    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(5), Inches(1.5),
                 "87.6%", font_size=80, color=GREEN, bold=True)

    add_text_box(slide, Inches(0.6), Inches(2.4), Inches(5), Inches(0.5),
                 "Overall Accuracy (Swin-Base, ImageNet pretrained)", font_size=16, color=TEXT_LIGHT)

    # Metrics grid
    metrics = [
        ("90.4%", "Balanced Accuracy"),
        ("89.2%", "F1 Macro"),
        ("0.800", "MCC"),
        ("0.977", "AUC-ROC"),
    ]
    x = Inches(0.6)
    for val, label in metrics:
        add_rounded_rect(slide, x, Inches(3.2), Inches(2.6), Inches(1.3),
                         RGBColor(0xED, 0xEB, 0xE5))
        add_text_box(slide, x, Inches(3.3), Inches(2.6), Inches(0.7),
                     val, font_size=28, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.9), Inches(2.6), Inches(0.4),
                     label, font_size=11, color=CITATION_C, alignment=PP_ALIGN.CENTER)
        x += Inches(2.8)

    # Per-class
    add_text_box(slide, Inches(0.6), Inches(4.8), Inches(5), Inches(0.4),
                 "Per-Class F1 Scores:", font_size=14, color=TEXT_LIGHT, bold=True)

    classes = [
        ("Mild", "0.867"), ("Moderate", "0.952"),
        ("Non-Demented", "0.896"), ("Very Mild", "0.851"),
    ]
    x = Inches(0.6)
    for cls, f1 in classes:
        add_text_box(slide, x, Inches(5.3), Inches(2.5), Inches(0.3),
                     f'{cls}: {f1}', font_size=13, color=TEXT_LIGHT)
        x += Inches(2.7)

    add_text_box(slide, Inches(0.6), Inches(5.8), Inches(10), Inches(0.4),
                 '100% recall on Moderate Demented (rarest class, n=12 in test set)',
                 font_size=14, color=GREEN, bold=True)

    # Comparison metrics chart from repo
    img = os.path.join(ALZ_DIR, "comparison_metrics.png")
    add_image_safe(slide, img, Inches(6.8), Inches(0.5), width=Inches(6.0))

    add_citation(slide,
                 "Author's experimental results, alzTheBatch repository  |  microsoft/swin-base-patch4-window7-224  |  +29pp over best CNN",
                 bg_dark=False)


def slide_18_cnn_vs_swin(prs, chart_path):
    """NEW: CNN vs Swin head-to-head comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "HEAD-TO-HEAD COMPARISON", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.7),
                 "CNN vs Swin Transformer: The Evidence",
                 font_size=28, color=TEXT_LIGHT, bold=True)

    # Side-by-side hero stats
    # CNN box
    add_rounded_rect(slide, Inches(0.6), Inches(1.8), Inches(3.5), Inches(2.0),
                     RGBColor(0xED, 0xEB, 0xE5), RED)
    add_text_box(slide, Inches(0.6), Inches(1.9), Inches(3.5), Inches(0.4),
                 "CNN (Combined Strategy)", font_size=13, color=RED, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(0.9),
                 "59.3%", font_size=42, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(3.5), Inches(0.4),
                 "Accuracy", font_size=12, color=CITATION_C, alignment=PP_ALIGN.CENTER)

    # Arrow
    add_text_box(slide, Inches(4.3), Inches(2.3), Inches(1.2), Inches(0.9),
                 '\u2192', font_size=36, color=COPPER, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.3), Inches(3.0), Inches(1.2), Inches(0.5),
                 '+29pp', font_size=14, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # Swin box
    add_rounded_rect(slide, Inches(5.6), Inches(1.8), Inches(3.5), Inches(2.0),
                     RGBColor(0xED, 0xEB, 0xE5), GREEN)
    add_text_box(slide, Inches(5.6), Inches(1.9), Inches(3.5), Inches(0.4),
                 "Swin Transformer", font_size=13, color=GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.6), Inches(2.3), Inches(3.5), Inches(0.9),
                 "87.6%", font_size=42, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.6), Inches(3.2), Inches(3.5), Inches(0.4),
                 "Accuracy", font_size=12, color=CITATION_C, alignment=PP_ALIGN.CENTER)

    # Chart
    add_image_safe(slide, chart_path, Inches(0.6), Inches(4.1), width=Inches(8.2))

    # Key insight box on right
    add_rounded_rect(slide, Inches(9.3), Inches(1.8), Inches(3.6), Inches(5.0),
                     RGBColor(0xED, 0xEB, 0xE5))
    add_text_box(slide, Inches(9.5), Inches(1.9), Inches(3.2), Inches(0.4),
                 "Key Factors", font_size=15, color=COPPER, bold=True)

    insights = [
        "Transfer learning from ImageNet provides rich visual features",
        "Shifted-window attention captures local + global patterns efficiently",
        "Differential learning rates preserve pretrained features",
        "Aggressive minority augmentation addresses class imbalance",
        "Validates thesis claims from Chapters 5\u20136",
    ]
    y = Inches(2.5)
    for ins in insights:
        add_text_box(slide, Inches(9.5), y, Inches(3.2), Inches(0.6),
                     f'\u2022  {ins}', font_size=11, color=TEXT_LIGHT)
        y += Inches(0.6)

    add_citation(slide,
                 "Author's experimental results  |  Liu et al., 2021 (Swin)  |  Deng et al., 2009 (ImageNet)",
                 bg_dark=False)


def slide_19_why_models_fail(prs, chart_path):
    """Why Models Fail (merged: data leakage + domain shift + shortcut learning)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "THE RECKONING", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.7),
                 "Why Models Fail in Practice", font_size=28, color=TEXT_LIGHT, bold=True)

    # Left: Data Leakage
    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(0.5),
                 "Data Leakage", font_size=18, color=RED, bold=True)

    add_image_safe(slide, chart_path, Inches(0.6), Inches(2.4), width=Inches(5.5))

    leakage_points = [
        "\u221228% accuracy when leakage is eliminated",
        "Only 4.5% of studies use proper methodology",
        "Same-subject slices in train AND test = fatal flaw",
    ]
    y = Inches(5.2)
    for pt in leakage_points:
        add_text_box(slide, Inches(0.6), y, Inches(5.8), Inches(0.35),
                     f'\u2022  {pt}', font_size=12, color=TEXT_LIGHT)
        y += Inches(0.38)

    # Divider
    add_shape_rect(slide, Inches(6.5), Inches(1.8), Pt(2), Inches(4.8), COPPER)

    # Right: Domain shift + Shortcut learning
    add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.5),
                 "Domain Shift & Shortcuts", font_size=18, color=COPPER, bold=True)

    right_points = [
        {'text': 'Domain Shift:', 'size': 14, 'color': TEXT_LIGHT, 'bold': True, 'spacing': 6},
        {'text': '\u2022  Models trained on ADNI drop to 71% on external data',
         'size': 12, 'color': TEXT_LIGHT, 'spacing': 4},
        {'text': '\u2022  Scanner differences, protocol variations, population demographics',
         'size': 12, 'color': TEXT_LIGHT, 'spacing': 4},
        {'text': '', 'size': 8, 'color': TEXT_LIGHT},
        {'text': 'Shortcut Learning:', 'size': 14, 'color': TEXT_LIGHT, 'bold': True, 'spacing': 10},
        {'text': '\u2022  Models exploit skull artifacts, not brain pathology',
         'size': 12, 'color': TEXT_LIGHT, 'spacing': 4},
        {'text': '\u2022  Background intensity patterns correlate with scanner/site',
         'size': 12, 'color': TEXT_LIGHT, 'spacing': 4},
        {'text': '\u2022  Grad-CAM reveals attention on non-brain regions',
         'size': 12, 'color': TEXT_LIGHT, 'spacing': 4},
        {'text': '', 'size': 8, 'color': TEXT_LIGHT},
        {'text': 'Both problems share a root cause:', 'size': 14, 'color': RED, 'bold': True, 'spacing': 10},
        {'text': 'Models learn spurious correlations rather than true disease biomarkers.',
         'size': 13, 'color': TEXT_LIGHT, 'spacing': 4},
    ]
    add_multiline_text(slide, Inches(7.0), Inches(2.3), Inches(5.8), Inches(4.2),
                       right_points)

    add_citation(slide, "Wen et al., 2020  |  Yagis et al., 2021  |  Geirhos et al., 2020", bg_dark=False)


def slide_20_accuracy_paradox(prs):
    """The Accuracy Paradox."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "THE ACCURACY PARADOX", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(6), Inches(0.7),
                 "When 95% Accuracy is Meaningless",
                 font_size=28, color=TEXT_LIGHT, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(1.7), Inches(2))

    points = [
        "Moderate Demented: only 1% of typical test sets",
        "A model predicting only 'Non-Demented' achieves >50% accuracy",
        "Balanced Accuracy, F1 Macro, and MCC are essential metrics",
        "Our CNN baseline: 55% accuracy but 0% recall on two classes",
        "Medical AI must be evaluated on the hardest cases, not the easiest",
    ]
    y = Inches(2.0)
    for pt in points:
        add_text_box(slide, Inches(0.6), y, Inches(6), Inches(0.5),
                     f'\u2022  {pt}', font_size=14, color=TEXT_LIGHT)
        y += Inches(0.55)

    # Image
    img = os.path.join(IMG_DIR, "limitations_class_imbalance.png")
    add_image_safe(slide, img, Inches(7.0), Inches(0.5), width=Inches(5.8))

    add_citation(slide, "Brodersen et al., 2010  |  Chicco & Jurman, 2020  |  Luque et al., 2019", bg_dark=False)


def slide_21_label_uncertainty(prs):
    """Label Uncertainty -- hero stat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "LABEL UNCERTAINTY", font_size=11, color=COPPER, bold=True)

    # Hero stat
    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5), Inches(1.5),
                 "71%", font_size=96, color=RED, bold=True)

    add_text_box(slide, Inches(0.6), Inches(3.0), Inches(5.5), Inches(0.5),
                 "of AD patients have mixed pathology at autopsy",
                 font_size=20, color=TEXT_LIGHT)

    add_accent_line(slide, Inches(0.6), Inches(3.7), Inches(2))

    points = [
        "Clinical labels used for training are probabilistic, not definitive",
        "Gold standard diagnosis requires post-mortem examination",
        "Co-pathology (TDP-43, Lewy bodies, vascular) creates noisy labels",
        "Irreducible error floor: even perfect models cannot exceed label accuracy",
        "Label noise affects both training signal and evaluation metrics",
    ]
    y = Inches(4.0)
    for pt in points:
        add_text_box(slide, Inches(0.6), y, Inches(12), Inches(0.4),
                     f'\u2022  {pt}', font_size=14, color=TEXT_LIGHT)
        y += Inches(0.5)

    add_citation(slide, "Kapasi et al., 2017  |  Beach et al., 2012  |  Schneider et al., 2007", bg_dark=False)


def slide_22_advances_road_ahead(prs):
    """Advances & The Road Ahead (merged)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "ADVANCES & THE ROAD AHEAD", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.7),
                 "Recent Progress & Future Directions",
                 font_size=28, color=TEXT_LIGHT, bold=True)

    # Left: Recent Advances
    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(0.5),
                 "Recent Advances", font_size=18, color=BLUE, bold=True)

    advances = [
        "Multimodal fusion (MRI + PET + genetics + CSF)",
        "Self-supervised pretraining for medical images",
        "Foundation models adapted to neuroimaging",
        "Transfer learning: ImageNet \u2192 medical tasks",
        "Federated learning for multi-site data",
    ]
    y = Inches(2.4)
    for adv in advances:
        add_text_box(slide, Inches(0.6), y, Inches(5.8), Inches(0.38),
                     f'\u2022  {adv}', font_size=13, color=TEXT_LIGHT)
        y += Inches(0.43)

    # Divider
    add_shape_rect(slide, Inches(6.5), Inches(1.8), Pt(2), Inches(4.8), COPPER)

    # Right: Future Directions
    add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.5),
                 "Three Future Pillars", font_size=18, color=COPPER, bold=True)

    pillars = [
        ("Standardized Benchmarks", "Unified evaluation protocols eliminating data leakage"),
        ("Clinical Interpretability", "Explanations that clinicians trust and understand"),
        ("Multi-Stage Diagnosis", "From binary AD/HC to the full continuum of decline"),
    ]
    y = Inches(2.4)
    for title, desc in pillars:
        add_rounded_rect(slide, Inches(7.0), y, Inches(5.5), Inches(1.1),
                         RGBColor(0xED, 0xEB, 0xE5))
        add_text_box(slide, Inches(7.2), y + Inches(0.1), Inches(5.1), Inches(0.4),
                     title, font_size=14, color=COPPER, bold=True)
        add_text_box(slide, Inches(7.2), y + Inches(0.5), Inches(5.1), Inches(0.5),
                     desc, font_size=12, color=TEXT_LIGHT)
        y += Inches(1.3)

    # Publications chart
    img = os.path.join(IMG_DIR, "Graph.png")
    add_image_safe(slide, img, Inches(0.6), Inches(4.8), width=Inches(5.8))

    add_citation(slide, "Zhang et al., 2022  |  Qui et al., 2022  |  Tanveer et al., 2024", bg_dark=False)


def slide_23_conclusion(prs):
    """Conclusion: The Methodological Triad (merged)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
                 "CONCLUSION", font_size=11, color=COPPER, bold=True)

    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(11), Inches(1.0),
                 "The Methodological Triad",
                 font_size=36, color=TEXT_DARK, bold=True)

    add_accent_line(slide, Inches(0.6), Inches(2.1), Inches(3))

    # Three non-negotiables
    triad = [
        ("01", "Rigorous Data Handling", "Subject-level splits, no data leakage, cross-site validation.\nWithout this, reported accuracy is meaningless.", BLUE),
        ("02", "Balanced Evaluation", "Balanced accuracy, F1 macro, MCC \u2014 not just accuracy.\nOur experiments confirm: baseline CNN = 55% accuracy, 0% minority recall.", COPPER),
        ("03", "Clinical Interpretability", "Grad-CAM, SHAP, clinical validation.\nModels must explain why, not just what.", GREEN),
    ]

    x = Inches(0.6)
    for num, title, desc, color in triad:
        add_rounded_rect(slide, x, Inches(2.6), Inches(3.8), Inches(2.8), DARK_ACCENT, color)
        add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(3.4), Inches(0.6),
                     num, font_size=32, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(3.3), Inches(3.4), Inches(0.5),
                     title, font_size=16, color=TEXT_DARK, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(3.8), Inches(3.4), Inches(1.4),
                     desc, font_size=12, color=CITATION_C)
        x += Inches(4.1)

    # Closing statement
    add_text_box(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.0),
                 '"Models are powerful. Data is insufficient. Trust is unearned."',
                 font_size=22, color=COPPER, bold=True, italic=True, alignment=PP_ALIGN.CENTER)

    add_citation(slide, "Wen et al., 2020  |  Author's analysis & experimental validation")


def slide_24_thank_you(prs):
    """Thank You / Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Logo
    logo = os.path.join(IMG_DIR, "logo_en.png")
    add_image_safe(slide, logo, Inches(0.6), Inches(0.4), width=Inches(2.8))

    add_accent_line(slide, Inches(0.6), Inches(2.8), Inches(4))

    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(12), Inches(1.0),
                 "Thank You", font_size=48, color=TEXT_DARK, bold=True)

    add_text_box(slide, Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
                 "Questions & Discussion", font_size=22, color=COPPER, italic=True)

    # Author info
    add_multiline_text(slide, Inches(0.6), Inches(5.4), Inches(6), Inches(1.5), [
        {'text': 'Paris Karageorgakis', 'size': 16, 'color': TEXT_DARK, 'bold': True},
        {'text': 'University of Piraeus, Department of Informatics', 'size': 13, 'color': CITATION_C, 'spacing': 6},
    ])

    # Links
    add_multiline_text(slide, Inches(7.0), Inches(5.4), Inches(5.5), Inches(1.5), [
        {'text': 'Thesis:  github.com/paris26/Thesis_Finale', 'size': 13, 'color': BLUE, 'spacing': 4},
        {'text': 'Experiments:  github.com/paris26/alzTheBatch', 'size': 13, 'color': BLUE, 'spacing': 6},
    ])


# ─────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("Building Thesis Presentation v2 (24 slides)")
    print("=" * 60)

    # Step 1: Generate charts
    print("\n[1/2] Generating matplotlib charts...")
    chart_prevalence = generate_prevalence_chart()
    print(f"  -> {chart_prevalence}")
    chart_ml = generate_classical_ml_chart()
    print(f"  -> {chart_ml}")
    chart_cnn_swin = generate_cnn_vs_swin_chart()
    print(f"  -> {chart_cnn_swin}")
    chart_leakage = generate_leakage_chart()
    print(f"  -> {chart_leakage}")

    # Step 2: Build presentation
    print("\n[2/2] Building presentation...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ACT I: THE PROBLEM (Slides 1-5)
    print("  ACT I: The Problem")
    slide_01_title(prs)
    print("    Slide 1: Title")
    slide_02_epidemic(prs, chart_prevalence)
    print("    Slide 2: The Silent Epidemic")
    slide_03_window_atn(prs)
    print("    Slide 3: 20-Year Window & AT(N)")
    slide_04_neuroimaging(prs)
    print("    Slide 4: Why Neuroimaging")
    slide_05_datasets(prs)
    print("    Slide 5: Benchmark Datasets")

    # ACT II: THE SCIENCE (Slides 6-18)
    print("  ACT II: The Science")
    slide_06_section_divider(prs)
    print("    Slide 6: Section Divider")
    slide_07_mri_fundamentals(prs)
    print("    Slide 7: MRI Fundamentals")
    slide_08_beyond_mri(prs)
    print("    Slide 8: Beyond MRI")
    slide_09_preprocessing(prs)
    print("    Slide 9: Preprocessing Pipeline")
    slide_10_signal_cleaning(prs)
    print("    Slide 10: Signal Cleaning")
    slide_11_skull_stripping(prs)
    print("    Slide 11: Skull Stripping")
    slide_12_vbm(prs)
    print("    Slide 12: VBM")
    slide_13_classical_ml(prs, chart_ml)
    print("    Slide 13: Classical ML")
    slide_14_cnns_to_transformers(prs)
    print("    Slide 14: CNNs to Transformers")
    slide_15_explainability(prs)
    print("    Slide 15: Explainability")
    slide_16_cnn_experiments(prs)
    print("    Slide 16: CNN Experiments [NEW]")
    slide_17_swin_results(prs)
    print("    Slide 17: Swin Results [NEW]")
    slide_18_cnn_vs_swin(prs, chart_cnn_swin)
    print("    Slide 18: CNN vs Swin [NEW]")

    # ACT III: THE RECKONING (Slides 19-23)
    print("  ACT III: The Reckoning")
    slide_19_why_models_fail(prs, chart_leakage)
    print("    Slide 19: Why Models Fail")
    slide_20_accuracy_paradox(prs)
    print("    Slide 20: Accuracy Paradox")
    slide_21_label_uncertainty(prs)
    print("    Slide 21: Label Uncertainty")
    slide_22_advances_road_ahead(prs)
    print("    Slide 22: Advances & Road Ahead")
    slide_23_conclusion(prs)
    print("    Slide 23: Conclusion")

    # BOOKEND (Slide 24)
    print("  BOOKEND")
    slide_24_thank_you(prs)
    print("    Slide 24: Thank You")

    # Save
    prs.save(OUTPUT)
    print(f"\n{'=' * 60}")
    print(f"SAVED: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"{'=' * 60}")


if __name__ == '__main__':
    main()
